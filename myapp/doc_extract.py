"""Extracts plain text from an uploaded PDF/DOCX/CSV/TXT so it can be fed to
the AI chat as normal text — no OCR, no external service, just the
document's own text layer. Scanned/image-only PDFs have no text layer to
extract and will raise ExtractError rather than return nothing useful.
"""
import csv
import io
import os
import re
from html.parser import HTMLParser

from docx import Document as DocxDocument
from pypdf import PdfReader

MAX_CHARS = 15_000     # extracted text is capped before it ever reaches the model
MAX_PDF_PAGES = 50
MAX_CSV_ROWS = 500


class _PlainTextHTMLParser(HTMLParser):
    """Dependency-free HTML-to-text fallback used when bs4 is unavailable."""

    _IGNORED_TAGS = {'script', 'style', 'noscript', 'template'}
    _BLOCK_TAGS = {
        'address', 'article', 'aside', 'blockquote', 'br', 'dd', 'div', 'dl',
        'dt', 'figcaption', 'figure', 'footer', 'h1', 'h2', 'h3', 'h4', 'h5',
        'h6', 'header', 'hr', 'li', 'main', 'nav', 'ol', 'p', 'pre', 'section',
        'table', 'tbody', 'td', 'tfoot', 'th', 'thead', 'tr', 'ul',
    }

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.parts = []
        self.ignored_depth = 0

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        if tag in self._IGNORED_TAGS:
            self.ignored_depth += 1
        elif not self.ignored_depth and tag in self._BLOCK_TAGS:
            self.parts.append('\n')

    def handle_startendtag(self, tag, attrs):
        if not self.ignored_depth and tag.lower() in self._BLOCK_TAGS:
            self.parts.append('\n')

    def handle_endtag(self, tag):
        tag = tag.lower()
        if self.ignored_depth:
            if tag in self._IGNORED_TAGS:
                self.ignored_depth -= 1
            return
        if tag in self._BLOCK_TAGS:
            self.parts.append('\n')

    def handle_data(self, data):
        if self.ignored_depth:
            return
        text = ' '.join(data.split())
        if text:
            self.parts.append(text)

    def text(self):
        text = ' '.join(self.parts)
        text = re.sub(r'[ \t]*\n[ \t]*', '\n', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()


class ExtractError(Exception):
    """Raised for any extraction failure with a message safe to show the user."""


def extract_pdf(file_bytes):
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
    except Exception as e:
        raise ExtractError(f"Could not read that PDF: {e}")

    if reader.is_encrypted:
        try:
            reader.decrypt('')
        except Exception:
            pass
        if reader.is_encrypted:
            raise ExtractError('That PDF is password-protected — remove the password and try again.')

    pages = []
    for page in reader.pages[:MAX_PDF_PAGES]:
        try:
            pages.append(page.extract_text() or '')
        except Exception:
            continue
    text = '\n\n'.join(p for p in pages if p.strip())
    if not text.strip():
        # Render scanned pages and run the same local OCR used for images.
        try:
            import base64
            import fitz
            from myapp.image_ocr import extract_data_uri
            scan = fitz.open(stream=file_bytes, filetype='pdf')
            ocr_pages = []
            for page in scan[:min(MAX_PDF_PAGES, 10)]:
                png = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).tobytes('png')
                uri = 'data:image/png;base64,' + base64.b64encode(png).decode('ascii')
                page_text = extract_data_uri(uri)
                if page_text:
                    ocr_pages.append(page_text)
            text = '\n\n'.join(ocr_pages)
        except Exception:
            text = ''
        if not text.strip():
            raise ExtractError("Couldn't find readable text in that PDF.")
    return text


def extract_docx(file_bytes):
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
    except Exception as e:
        raise ExtractError(f"Could not read that document: {e}")

    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(' | '.join(cell.text for cell in row.cells))
    text = '\n'.join(parts)
    if not text.strip():
        raise ExtractError('That document appears to be empty.')
    return text


def extract_csv(file_bytes):
    try:
        text_data = file_bytes.decode('utf-8-sig', errors='replace')
        rows = list(csv.reader(io.StringIO(text_data)))
    except Exception as e:
        raise ExtractError(f"Could not read that CSV: {e}")
    if not rows:
        raise ExtractError('That CSV appears to be empty.')
    truncated_rows = rows[:MAX_CSV_ROWS]
    lines = [', '.join(cell.strip() for cell in row) for row in truncated_rows]
    if len(rows) > MAX_CSV_ROWS:
        lines.append(f'... ({len(rows) - MAX_CSV_ROWS} more rows not shown)')
    try:
        import pandas as pd
        frame = pd.read_csv(io.BytesIO(file_bytes))
        summary = frame.describe(include='all').fillna('').to_string()
        lines.extend(['', 'Data summary:', summary])
    except Exception:
        pass
    return '\n'.join(lines)


def extract_txt(file_bytes):
    try:
        text = file_bytes.decode('utf-8-sig', errors='replace')
    except Exception as e:
        raise ExtractError(f"Could not read that file: {e}")
    if not text.strip():
        raise ExtractError('That file appears to be empty.')
    return text


def extract_xlsx(file_bytes):
    try:
        from openpyxl import load_workbook
        book = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines = []
        for sheet in book.worksheets:
            lines.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                lines.append(' | '.join('' if value is None else str(value) for value in row))
                if len(lines) >= MAX_CSV_ROWS:
                    lines.append('... (remaining spreadsheet rows omitted)')
                    return '\n'.join(lines)
        return '\n'.join(lines)
    except Exception as e:
        raise ExtractError(f"Could not read that spreadsheet: {e}")


def extract_pptx(file_bytes):
    try:
        from pptx import Presentation
        deck = Presentation(io.BytesIO(file_bytes))
        lines = []
        for number, slide in enumerate(deck.slides, 1):
            parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
            if parts:
                lines.append(f"Slide {number}:\n" + '\n'.join(parts))
        return '\n\n'.join(lines)
    except Exception as e:
        raise ExtractError(f"Could not read that presentation: {e}")


def extract_html(file_bytes):
    # BeautifulSoup gives the best results when installed, but HTML upload is
    # useful enough that it should not fail solely because an optional package
    # is missing from a partially provisioned environment. Python's standard
    # HTMLParser provides a safe local fallback with no network or subprocess.
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes, 'html.parser')
        for element in soup(['script', 'style', 'noscript', 'template']):
            element.decompose()
        text = '\n'.join(s.strip() for s in soup.stripped_strings if s.strip())
    except (ImportError, ModuleNotFoundError):
        try:
            parser = _PlainTextHTMLParser()
            parser.feed(file_bytes.decode('utf-8-sig', errors='replace'))
            parser.close()
            text = parser.text()
        except Exception as e:
            raise ExtractError(f"Could not read that HTML file: {e}")
    except Exception as e:
        raise ExtractError(f"Could not read that HTML file: {e}")

    if not text.strip():
        raise ExtractError('That HTML file does not contain readable text.')
    return text


_EXTRACTORS = {
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.csv': extract_csv,
    '.txt': extract_txt,
    '.xlsx': extract_xlsx,
    '.xlsm': extract_xlsx,
    '.pptx': extract_pptx,
    '.html': extract_html,
    '.htm': extract_html,
}

# Source-code and config formats are plain text, so the same safe decoder used
# for TXT works for them. This also lets "Start coding" operate on the actual
# source rather than forcing users to rename a file to .txt first.
CODE_EXTENSIONS = (
    '.css', '.js', '.mjs', '.cjs', '.ts', '.tsx', '.jsx', '.py', '.java',
    '.c', '.cpp', '.h', '.hpp', '.cs', '.php', '.rb', '.go', '.rs', '.swift',
    '.kt', '.kts', '.sql', '.json', '.xml', '.yaml', '.yml', '.md', '.vue',
    '.svelte', '.sh', '.ps1', '.toml', '.ini',
)
_EXTRACTORS.update({extension: extract_txt for extension in CODE_EXTENSIONS})

# These formats can be returned to coding mode byte-for-byte (after decoding),
# unlike PDF/DOCX/XLSX/PPTX where only their extracted textual content is
# available and binary formatting cannot be reconstructed in chat.
EDITABLE_SOURCE_EXTENSIONS = frozenset(CODE_EXTENSIONS + ('.html', '.htm', '.txt', '.csv'))

SUPPORTED_EXTENSIONS = tuple(_EXTRACTORS.keys())


def extract(filename, file_bytes):
    """Returns (text, truncated). Raises ExtractError on any failure."""
    ext = os.path.splitext(filename or '')[1].lower()
    extractor = _EXTRACTORS.get(ext)
    if not extractor:
        raise ExtractError(
            f"Unsupported file type '{ext or filename}'. Supported: documents, spreadsheets, presentations, HTML, and common source-code files."
        )
    text = extractor(file_bytes)
    truncated = len(text) > MAX_CHARS
    if truncated:
        text = text[:MAX_CHARS]
    return text, truncated


def extract_editable_source(filename, file_bytes, extracted_text, extracted_truncated=False):
    """Return the best text to edit in coding mode plus its truncation flag.

    HTML/details extraction intentionally strips tags and scripts, which is
    correct for summarising but wrong for editing. For source-backed formats
    preserve the original source. Binary office formats fall back to their
    extracted text so they can still use the same editing workflow.
    """
    ext = os.path.splitext(filename or '')[1].lower()
    if ext in EDITABLE_SOURCE_EXTENSIONS:
        source = file_bytes.decode('utf-8-sig', errors='replace')
        truncated = len(source) > MAX_CHARS
    else:
        source = extracted_text
        truncated = bool(extracted_truncated) or len(source) > MAX_CHARS
    return source[:MAX_CHARS] if truncated else source, truncated
