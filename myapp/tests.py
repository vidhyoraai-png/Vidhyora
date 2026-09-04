import base64
import datetime
import io
import json
import tempfile
from datetime import timedelta
from decimal import Decimal
from types import SimpleNamespace
from unittest.mock import Mock, patch
from zoneinfo import ZoneInfo

from django.contrib.auth.models import User
from django.core.cache import cache
from django.core.files.uploadedfile import SimpleUploadedFile
from django.db import OperationalError
from django.http import HttpResponse
from django.contrib.sessions.models import Session
from django.test import RequestFactory
from django.test import TestCase, override_settings
from django.utils import timezone
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from . import (
    ai_chat, business_info, company_knowledge, doc_extract, dropbox_backup,
    dropbox_images, file_convert, image_generation, privacy, request_router,
    web_search,
)
from .middleware import CanonicalHostMiddleware, PublicAssetCacheMiddleware
from .models import ActiveUserSession, AIGeneratedFile, AIBlock, AIConversation, AIMessage, AINote, AIReport, AIUserImage, GitHubConnection, PWASettings, StoreProfile
from .views import (
    AI_CURRENT_CONVERSATION_SESSION_KEY, _ai_document_instruction,
    _ai_excel_bytes, _ai_generated_file_spec, _ai_pdf_bytes,
    _ai_powerpoint_bytes, _ai_word_document_bytes,
    _extract_ai_generated_file_content, _strip_fake_download_links,
    site_customization_context,
)


class AIConversationPersistenceTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username='chat-persistence@example.com',
            email='chat-persistence@example.com',
            password='test-password-123',
        )
        StoreProfile.objects.create(user=self.user, phone='9999999999')

    def test_open_conversation_is_restored_on_authenticated_refresh(self):
        self.client.force_login(self.user)
        older = AIConversation.objects.create(user=self.user, title='Older chat')
        selected = AIConversation.objects.create(user=self.user, title='Selected chat')
        AIMessage.objects.create(conversation=selected, role=AIMessage.ROLE_USER, content='Keep this open')

        response = self.client.get(f'/AI/api/conversations/{selected.id}/')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(self.client.session[AI_CURRENT_CONVERSATION_SESSION_KEY], selected.id)

        response = self.client.get('/AI/')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.context['ai_resume_conversation_id'], selected.id)
        self.assertContains(response, f'var AI_RESUME_CONVERSATION_ID = {selected.id};')
        self.assertNotEqual(older.id, response.context['ai_resume_conversation_id'])

    def test_refresh_falls_back_to_newest_owned_conversation(self):
        self.client.force_login(self.user)
        conversation = AIConversation.objects.create(user=self.user, title='Latest chat')

        response = self.client.get('/AI/')

        self.assertEqual(response.context['ai_resume_conversation_id'], conversation.id)
        self.assertEqual(self.client.session[AI_CURRENT_CONVERSATION_SESSION_KEY], conversation.id)

    def test_root_url_serves_the_ai_homepage(self):
        response = self.client.get('/')

        self.assertEqual(response.status_code, 200)
        self.assertTemplateUsed(response, 'ai.html')
        self.assertEqual(response.context['ai_default_model'], 'quick')

    def test_flux_model_is_available_in_the_ai_picker(self):
        self.assertIn(ai_chat.FLUX_KLEIN_4B_MODEL_KEY, ai_chat.MODELS)
        self.assertEqual(ai_chat.MODELS[ai_chat.FLUX_KLEIN_4B_MODEL_KEY]['label'], 'FLUX.2 Klein 4B')

        response = self.client.get('/AI/')
        self.assertContains(response, 'FLUX.2 Klein 4B')

    def test_guest_chat_survives_login_and_remains_selected(self):
        session = self.client.session
        session.save()
        guest_session_key = session.session_key
        conversation = AIConversation.objects.create(
            session_key=guest_session_key,
            title='Guest chat',
        )
        AIMessage.objects.create(conversation=conversation, role=AIMessage.ROLE_USER, content='Before login')
        session[AI_CURRENT_CONVERSATION_SESSION_KEY] = conversation.id
        session.save()

        response = self.client.post(
            '/AI/api/login/',
            data=json.dumps({
                'identifier': self.user.email,
                'password': 'test-password-123',
            }),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        conversation.refresh_from_db()
        self.assertEqual(conversation.user, self.user)
        self.assertEqual(conversation.session_key, '')
        self.assertEqual(self.client.session[AI_CURRENT_CONVERSATION_SESSION_KEY], conversation.id)

        response = self.client.get('/AI/')
        self.assertEqual(response.context['ai_resume_conversation_id'], conversation.id)
        response = self.client.get(f'/AI/api/conversations/{conversation.id}/')
        self.assertEqual(response.json()['messages'][0]['content'], 'Before login')


class SingleDeviceLoginTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username='one-device@example.com', email='one-device@example.com',
            password='test-password-123',
        )
        StoreProfile.objects.create(user=self.user, phone='9777777777')

    def _login(self, client):
        return client.post(
            '/AI/api/login/',
            data=json.dumps({
                'identifier': self.user.email,
                'password': 'test-password-123',
            }),
            content_type='application/json',
        )

    def test_new_login_immediately_invalidates_the_previous_device(self):
        first_device = self.client_class()
        second_device = self.client_class()

        self.assertEqual(self._login(first_device).status_code, 200)
        first_key = first_device.session.session_key
        self.assertEqual(
            ActiveUserSession.objects.get(user=self.user).session_key,
            first_key,
        )

        self.assertEqual(self._login(second_device).status_code, 200)
        second_key = second_device.session.session_key

        self.assertNotEqual(first_key, second_key)
        self.assertFalse(Session.objects.filter(session_key=first_key).exists())
        self.assertEqual(
            ActiveUserSession.objects.get(user=self.user).session_key,
            second_key,
        )
        self.assertEqual(first_device.get('/AI/api/account/').status_code, 401)
        self.assertEqual(second_device.get('/AI/api/account/').status_code, 200)

    def test_logging_out_releases_the_device_slot(self):
        self.assertEqual(self._login(self.client).status_code, 200)
        self.assertTrue(ActiveUserSession.objects.filter(user=self.user).exists())

        response = self.client.post('/AI/api/logout/')

        self.assertEqual(response.status_code, 200)
        self.assertFalse(ActiveUserSession.objects.filter(user=self.user).exists())

    def test_django_admin_login_uses_the_same_single_device_slot(self):
        self.user.is_staff = True
        self.user.is_superuser = True
        self.user.save(update_fields=['is_staff', 'is_superuser'])
        ai_device = self.client_class()
        admin_device = self.client_class()
        self.assertEqual(self._login(ai_device).status_code, 200)
        old_key = ai_device.session.session_key

        response = admin_device.post('/admin/login/?next=/admin/', {
            'username': self.user.username,
            'password': 'test-password-123',
            'next': '/admin/',
        })

        self.assertEqual(response.status_code, 302)
        self.assertFalse(Session.objects.filter(session_key=old_key).exists())
        self.assertEqual(
            ActiveUserSession.objects.get(user=self.user).session_key,
            admin_device.session.session_key,
        )
        self.assertEqual(ai_device.get('/AI/api/account/').status_code, 401)


class NVIDIAImageGenerationTests(TestCase):
    PNG_BYTES = b'\x89PNG\r\n\x1a\nmock-image'

    @override_settings(
        NVIDIA_FLUX_API_KEY='test-flux-key',
        NVIDIA_FLUX_EDIT_API_KEY='test-edit-key',
    )
    @patch('myapp.image_generation.requests.post')
    def test_text_prompt_uses_flux_generation_schema(self, post):
        post.return_value.status_code = 200
        post.return_value.json.return_value = {
            'artifacts': [{'base64': base64.b64encode(self.PNG_BYTES).decode()}],
        }

        result = image_generation.generate_image('A futuristic learning robot')

        self.assertEqual(result.content, self.PNG_BYTES)
        self.assertEqual(result.extension, 'png')
        request = post.call_args
        self.assertNotIn('mode', request.kwargs['json'])
        self.assertEqual(request.kwargs['json']['steps'], 4)
        self.assertNotIn('image', request.kwargs['json'])
        self.assertEqual(request.kwargs['headers']['Authorization'], 'Bearer test-flux-key')

    @override_settings(
        NVIDIA_FLUX_API_KEY='test-flux-key',
        NVIDIA_FLUX_EDIT_API_KEY='test-edit-key',
    )
    @patch('myapp.image_generation.requests.post')
    def test_attached_image_uses_flux_editing_schema(self, post):
        post.return_value.status_code = 200
        post.return_value.json.return_value = {
            'artifacts': [{'base64': base64.b64encode(self.PNG_BYTES).decode()}],
        }
        source_bytes = io.BytesIO()
        Image.new('RGBA', (8, 6), (0, 0, 255, 128)).save(source_bytes, format='PNG')
        source = 'data:image/png;base64,' + base64.b64encode(source_bytes.getvalue()).decode()

        image_generation.generate_image('Remove the background', source)

        body = post.call_args.kwargs['json']
        self.assertNotIn('mode', body)
        self.assertEqual(len(body['image']), 1)
        self.assertTrue(body['image'][0].startswith('data:image/jpeg;base64,'))
        self.assertEqual(post.call_args.kwargs['headers']['Authorization'], 'Bearer test-edit-key')

    @override_settings(NVIDIA_FLUX_API_KEY='test-flux-key')
    @patch('myapp.image_generation.requests.post')
    def test_upstream_rate_limit_is_safe_and_actionable(self, post):
        post.return_value.status_code = 429

        with self.assertRaises(image_generation.ImageGenerationError) as raised:
            image_generation.generate_image('Create a poster')

        self.assertEqual(raised.exception.status_code, 429)
        self.assertIn('limit', str(raised.exception).lower())

    @override_settings(NVIDIA_FLUX_API_KEY='')
    @patch('myapp.image_generation.requests.post')
    def test_missing_flux_key_fails_without_an_upstream_request(self, post):
        with self.assertRaises(image_generation.ImageGenerationError) as raised:
            image_generation.generate_image('Create a poster')

        self.assertIn('NVIDIA_FLUX_API_KEY', str(raised.exception))
        post.assert_not_called()

    def setUp(self):
        cache.clear()
        self.user = User.objects.create_user(
            username='flux-tests@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(self.user)

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/result.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/result.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_flux_chat_turn_saves_and_returns_generated_image(self, generate, save, storage_url):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')

        response = self.client.post(
            '/AI/api/send/',
            data=json.dumps({
                'message': 'A futuristic EduTrellis AI robot',
                'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
            }),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.content.decode(), '')
        self.assertEqual(response['X-Generated-Image-Url'], '/media/ai_generated/result.png')
        self.assertEqual(response['X-Request-Category'], 'image_generation')
        generate.assert_called_once_with('A futuristic EduTrellis AI robot', None)
        assistant = AIMessage.objects.get(role=AIMessage.ROLE_ASSISTANT)
        self.assertEqual(assistant.model_key, ai_chat.FLUX_KLEIN_4B_MODEL_KEY)
        self.assertEqual(assistant.image_data, '/media/ai_generated/result.png')

        history = self.client.get(f'/AI/api/conversations/{assistant.conversation_id}/').json()
        self.assertEqual(history['messages'][-1]['image_data'], '/media/ai_generated/result.png')

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/edited.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/edited.png')
    @patch('myapp.views.image_generation.generate_image')
    @patch('myapp.views.image_ocr.extract_data_uri', return_value='')
    def test_flux_chat_turn_edits_attachment_instead_of_routing_to_vision(
        self, ocr, generate, save, storage_url,
    ):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')
        source = 'data:image/png;base64,AA=='

        response = self.client.post(
            '/AI/api/send/',
            data=json.dumps({
                'message': 'Make the background blue',
                'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
                'image': source,
            }),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.content.decode(), '')
        self.assertEqual(response['X-Routed-Model-Key'], ai_chat.FLUX_KLEIN_4B_MODEL_KEY)
        self.assertEqual(response['X-Request-Category'], 'image_edit')
        generate.assert_called_once_with('Make the background blue', source)

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/chatgpt.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/chatgpt.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_chatgpt_image_generation_hides_the_flux_worker_label(self, generate, save, storage_url):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')

        response = self.client.post(
            '/AI/api/send/',
            data=json.dumps({
                'message': 'Generate an image of an orange robot',
                'model': ai_chat.CHATGPT_56_MODEL_KEY,
            }),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['X-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response['X-Routed-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        assistant = AIMessage.objects.get(role=AIMessage.ROLE_ASSISTANT)
        self.assertEqual(assistant.model_key, ai_chat.CHATGPT_56_MODEL_KEY)

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/cat.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/cat.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_report_9_hinglish_bna_prompt_routes_to_image_generation(self, generate, save, storage_url):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')
        prompt = 'Cat ka image bna ke do'

        response = self.client.post(
            '/AI/api/send/',
            data=json.dumps({'message': prompt, 'model': ai_chat.CHATGPT_56_MODEL_KEY}),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['X-Routed-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response['X-Request-Category'], 'image_generation')
        generate.assert_called_once_with(prompt, None)

    @patch('myapp.views.image_generation.generate_image')
    def test_chatgpt_image_errors_never_expose_internal_provider_or_model_names(self, generate):
        generate.side_effect = image_generation.ImageGenerationError(
            "That image request was blocked by NVIDIA's content filter. Try a different prompt or image.",
            status_code=400,
        )

        response = self.client.post(
            '/AI/api/send/',
            data=json.dumps({
                'message': 'Generate an image of a person sitting on a bus',
                'model': ai_chat.CHATGPT_56_MODEL_KEY,
            }),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response['X-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response['X-Routed-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        detail = response.json()['detail']
        self.assertEqual(
            detail,
            'That image request was blocked by the safety filter. Try a different prompt or image.',
        )
        for hidden_name in ('NVIDIA', 'FLUX', 'Nemotron', 'Black Forest'):
            self.assertNotIn(hidden_name.lower(), detail.lower())


class ImageAspectRatioFromPromptTests(TestCase):
    """Reports #45, #47 and #48: every image came back a 1024x1024 square.

    The sizes asserted here were confirmed against the live FLUX endpoint —
    it returns exactly the requested dimensions, and rejects anything over
    1,062,400 pixels.
    """

    def test_every_supported_size_is_within_the_api_pixel_ceiling(self):
        for ratio, (width, height) in image_generation._ASPECT_SIZES.items():
            with self.subTest(ratio=ratio):
                self.assertLessEqual(width * height, image_generation.MAX_PIXELS)
                # Diffusion models want dimensions on a 32px grid.
                self.assertEqual((width % 32, height % 32), (0, 0))

    def test_wallpaper_request_becomes_landscape(self):
        # Report #48: "wallpaper 4k resolution" returned a square.
        self.assertEqual(
            image_generation.resolve_dimensions(
                'Generate image of spiderman with black background wallpaper 4k resolution'
            ),
            (1344, 768),
        )

    def test_size_written_with_a_dot_is_read_as_a_ratio(self):
        # Report #47: "size 9.12" means 9:12, i.e. a 3:4 portrait.
        self.assertEqual(
            image_generation.resolve_dimensions(
                'Create image good morning size 9.12 with motivational msg'
            ),
            (896, 1152),
        )

    def test_instagram_post_becomes_a_feed_shaped_portrait(self):
        # Report #45: "turn in to instagram post".
        self.assertEqual(
            image_generation.resolve_dimensions('turn in to instagram post'), (896, 1120),
        )

    def test_orientation_words_pick_the_matching_shape(self):
        cases = {
            'make a youtube thumbnail of a cat': (1344, 768),
            'instagram story for diwali sale': (768, 1344),
            'a poster for our new shop': (896, 1152),
            'profile picture of a lion': (1024, 1024),
            'landscape photo of a beach': (1344, 768),
            'portrait of a woman': (768, 1344),
            '1920x1080 image of a sunset': (1344, 768),
            'image with 4:5 ratio': (896, 1120),
        }
        for prompt, expected in cases.items():
            with self.subTest(prompt=prompt):
                self.assertEqual(image_generation.resolve_dimensions(prompt), expected)

    def test_a_phone_wallpaper_beats_the_plain_wallpaper_cue(self):
        # Both words match; the more specific one has to win.
        self.assertEqual(
            image_generation.resolve_dimensions('phone wallpaper of mountains'), (768, 1344),
        )
        self.assertEqual(
            image_generation.resolve_dimensions('desktop wallpaper of mountains'), (1344, 768),
        )

    def test_numbers_that_are_not_aspect_ratios_leave_the_default_alone(self):
        for prompt in (
            'draw a cat',
            'ChatGPT 5.6 logo',                 # a version number, not 5:6
            'good morning image at 9:30 am',    # a clock time, not 9:30
            'image of a 16 year old birthday cake',
            'generate an image of the number 1/2',
        ):
            with self.subTest(prompt=prompt):
                self.assertEqual(
                    image_generation.resolve_dimensions(prompt),
                    image_generation.DEFAULT_SIZE,
                )

    @patch('myapp.image_generation.requests.post')
    def test_the_resolved_size_is_what_actually_reaches_the_api(self, post):
        post.return_value = Mock(
            status_code=200,
            json=Mock(return_value={'artifacts': [{
                'base64': base64.b64encode(
                    base64.b64decode(
                        b'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQ'
                        b'DwAEhQGAhKmMIQAAAABJRU5ErkJggg=='
                    )
                ).decode(),
            }]}),
        )

        with override_settings(NVIDIA_FLUX_API_KEY='test-key'):
            image_generation.generate_image('a 16:9 banner for my website')

        body = post.call_args.kwargs['json']
        self.assertEqual((body['width'], body['height']), (1344, 768))


class ImageGallerySurvivesChatDeletionTests(TestCase):
    """Deleting a chat must not destroy the images generated inside it."""

    PNG_BYTES = base64.b64decode(
        b'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=='
    )

    def setUp(self):
        cache.clear()
        self.user = User.objects.create_user(
            username='gallery-user', email='gallery@example.com',
            password='test-password-123', is_staff=True,
        )
        StoreProfile.objects.get_or_create(user=self.user)
        self.client.force_login(self.user)

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/kept.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/kept.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_image_stays_in_the_gallery_after_its_chat_is_deleted(
        self, generate, save, storage_url,
    ):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')
        send = self.client.post(
            '/AI/api/send/',
            data=json.dumps({
                'message': 'a calm blue lake',
                'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
            }),
            content_type='application/json',
        )
        self.assertEqual(send.status_code, 200)
        conversation_id = int(send['X-Conversation-Id'])

        gallery = self.client.get('/AI/api/account/').json()['images']
        self.assertEqual([i['url'] for i in gallery], ['/media/ai_generated/kept.png'])

        delete = self.client.post(f'/AI/api/conversations/{conversation_id}/delete/')
        self.assertEqual(delete.status_code, 200)
        self.assertEqual(AIMessage.objects.count(), 0)

        # The whole point: the picture is still there.
        after = self.client.get('/AI/api/account/').json()['images']
        self.assertEqual([i['url'] for i in after], ['/media/ai_generated/kept.png'])
        self.assertEqual(after[0]['prompt'], 'a calm blue lake')
        # ...and it no longer points at a chat that does not exist.
        self.assertIsNone(after[0]['conversation_id'])

    def test_the_gallery_never_shows_another_account_images(self):
        other = User.objects.create_user('other-user', password='pw')
        AIUserImage.objects.create(user=other, url='/media/ai_generated/theirs.png')
        AIUserImage.objects.create(user=self.user, url='/media/ai_generated/mine.png')

        images = self.client.get('/AI/api/account/').json()['images']
        self.assertEqual([i['url'] for i in images], ['/media/ai_generated/mine.png'])

    def test_there_is_no_endpoint_that_deletes_a_gallery_image(self):
        # "cannot be deleted": nothing in the app removes these rows, so a
        # stray URL should not quietly become a delete route.
        image = AIUserImage.objects.create(user=self.user, url='/media/ai_generated/x.png')
        for path in (
            f'/AI/api/images/{image.pk}/delete/',
            f'/AI/api/account/images/{image.pk}/delete/',
        ):
            with self.subTest(path=path):
                self.assertEqual(self.client.post(path).status_code, 404)
        self.assertTrue(AIUserImage.objects.filter(pk=image.pk).exists())


class DashboardReportStatsTests(TestCase):
    """Solved/unresolved totals, and tiles that filter when clicked."""

    def setUp(self):
        cache.clear()
        staff = User.objects.create_user('report-staff', password='pw', is_staff=True)
        StoreProfile.objects.create(user=staff)
        self.client.force_login(staff)
        conversation = AIConversation.objects.create(title='c')
        for index in range(3):
            AIReport.objects.create(
                conversation=conversation, user_prompt=f'p{index}',
                reported_reply='r', explanation=f'wrong {index}',
                status=AIReport.STATUS_OPEN,
            )
        for index in range(2):
            AIReport.objects.create(
                conversation=conversation, user_prompt=f'q{index}',
                reported_reply='r', explanation=f'fixed {index}',
                status=AIReport.STATUS_RESOLVED,
            )

    def test_counts_of_solved_and_unresolved_are_shown(self):
        stats = self.client.get('/store/dashboard/ai/reports/').context['report_stats']
        self.assertEqual(stats['total'], 5)
        self.assertEqual(stats['open'], 3)
        self.assertEqual(stats['resolved'], 2)
        self.assertEqual(stats['resolved_percent'], 40)

    def test_status_filter_narrows_the_listing_but_not_the_totals(self):
        response = self.client.get('/store/dashboard/ai/reports/?status=resolved')
        self.assertEqual(len(response.context['reports']), 2)
        # Totals must stay whole so the other tiles don't read zero.
        self.assertEqual(response.context['report_stats']['open'], 3)
        self.assertEqual(response.context['report_stats']['total'], 5)

    def test_an_unknown_status_filter_is_ignored_rather_than_erroring(self):
        response = self.client.get('/store/dashboard/ai/reports/?status=banana')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.context['status_filter'], '')
        self.assertEqual(len(response.context['reports']), 5)


class DashboardSignupFilterTests(TestCase):
    """Clicking a signup stat tile lists exactly those accounts."""

    def setUp(self):
        cache.clear()
        staff = User.objects.create_user('signup-staff', password='pw', is_staff=True)
        StoreProfile.objects.create(user=staff)
        self.client.force_login(staff)

        located = User.objects.create_user('located-user', password='pw')
        StoreProfile.objects.create(
            user=located, location_consent=StoreProfile.LOCATION_GRANTED,
        )
        payer = User.objects.create_user('paying-user', password='pw')
        StoreProfile.objects.create(user=payer, manual_amount_paid=Decimal('499.00'))
        plain = User.objects.create_user('plain-user', password='pw')
        StoreProfile.objects.create(user=plain)

    def test_location_tile_lists_only_accounts_with_location_enabled(self):
        response = self.client.get('/store/dashboard/signups/?filter=location')
        self.assertEqual(
            [u.username for u in response.context['users']], ['located-user'],
        )
        self.assertEqual(response.context['filtered_count'], 1)
        self.assertIn('location enabled', response.context['filter_label'])

    def test_paid_tile_lists_only_accounts_with_a_recorded_payment(self):
        response = self.client.get('/store/dashboard/signups/?filter=paid')
        self.assertEqual([u.username for u in response.context['users']], ['paying-user'])

    def test_no_filter_lists_everyone(self):
        response = self.client.get('/store/dashboard/signups/')
        self.assertEqual(len(response.context['users']), 4)
        self.assertEqual(response.context['active_filter'], '')

    def test_an_unknown_filter_is_ignored_rather_than_erroring(self):
        response = self.client.get('/store/dashboard/signups/?filter=banana')
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.context['active_filter'], '')
        self.assertEqual(len(response.context['users']), 4)

    def test_the_filter_combines_with_the_search_box(self):
        response = self.client.get('/store/dashboard/signups/?filter=location&q=located')
        self.assertEqual([u.username for u in response.context['users']], ['located-user'])


class GeneratedFileQualityTests(TestCase):
    """The 'create a pdf/doc for me' failures seen in real use.

    Three separate bugs showed up in one screenshot set: two download links
    under one reply, a PDF whose only text was a fabricated link, and 'create
    doc file for me' producing a .txt.
    """

    def test_doc_and_word_requests_produce_a_real_word_file(self):
        cases = {
            'create doc file for me': 'generated.docx',
            'make a summarise note in word file': 'generated.docx',
            'create a word document about our pricing': 'generated.docx',
            'create this data in pdf file and give to me': 'generated.pdf',
            'give me an excel sheet of this': 'generated.xlsx',
            'make a presentation on solar energy': 'generated.pptx',
        }
        for prompt, expected in cases.items():
            with self.subTest(prompt=prompt):
                self.assertEqual(_ai_generated_file_spec(prompt)['file_name'], expected)

    def test_a_model_invented_download_link_is_removed_from_the_reply(self):
        # Verbatim from the failing screenshot, fake token and all.
        reply = ('[Download generated.pdf](http://127.0.0.1:8000/AI/api/files/'
                 '12345678-1234-1234-1234-123456789012/download/)')
        self.assertEqual(_strip_fake_download_links(reply), '')

    def test_a_reply_that_is_only_a_fake_link_creates_no_file(self):
        # This was becoming a PDF whose entire contents were the link text.
        reply = ('Here you go!\n\n[Download generated.pdf](http://x/AI/api/files/'
                 '12345678-1234-1234-1234-123456789012/download/)')
        self.assertNotIn('AI/api/files', _strip_fake_download_links(reply))
        self.assertEqual(_extract_ai_generated_file_content(reply), 'Here you go!')

    def test_a_clarifying_question_creates_no_file(self):
        # This was handed back as a .txt containing the question itself.
        for reply in (
            "I can help you create a document file. What topic or content would you "
            "like the document to cover?",
            "Please let me know what you'd like the document to contain.",
            "Could you specify the details you want included?",
        ):
            with self.subTest(reply=reply[:40]):
                self.assertEqual(_extract_ai_generated_file_content(reply), '')

    def test_real_content_still_reaches_the_file(self):
        reply = "Here it is.\n\n```markdown\n# Title\n\nSome **real** content.\n```"
        self.assertEqual(
            _extract_ai_generated_file_content(reply),
            '# Title\n\nSome **real** content.',
        )

    def test_generated_pdf_contains_the_actual_text(self):
        content = '# Quarterly Report\n\nRevenue grew by **18%**.\n\n- Delhi: 42\n- Mumbai: 31'
        reader = PdfReader(io.BytesIO(_ai_pdf_bytes(content)))
        text = '\n'.join(page.extract_text() or '' for page in reader.pages)

        self.assertIn('Quarterly Report', text)
        self.assertIn('Revenue grew by 18%', text)
        self.assertIn('Delhi: 42', text)
        # Markdown markers must not survive into a finished document.
        self.assertNotIn('**', text)

    def test_generated_word_file_uses_real_headings_and_bold(self):
        content = '# Title\n\nRevenue grew by **18%** this quarter.\n\n- One\n- Two\n\n1. First'
        document = Document(io.BytesIO(_ai_word_document_bytes(content)))
        styles = [p.style.name for p in document.paragraphs if p.text.strip()]
        texts = [p.text for p in document.paragraphs if p.text.strip()]

        self.assertIn('Heading 1', styles)
        self.assertIn('List Bullet', styles)
        self.assertIn('List Number', styles)
        self.assertIn('Revenue grew by 18% this quarter.', texts)
        for line in texts:
            self.assertNotIn('**', line)
        # "18%" must be a genuine bold run, not asterisks in the text.
        body = next(p for p in document.paragraphs if p.text.startswith('Revenue'))
        self.assertTrue(any(run.bold and '18%' in run.text for run in body.runs))

    def test_markdown_tables_and_links_render_as_readable_text(self):
        content = ('| Region | Clients |\n| --- | --- |\n| North | 42 |\n\n'
                   'See [our site](https://example.com) for more.')
        text = '\n'.join(
            page.extract_text() or ''
            for page in PdfReader(io.BytesIO(_ai_pdf_bytes(content))).pages
        )

        self.assertIn('Region', text)
        self.assertIn('North', text)
        self.assertNotIn('---', text)
        self.assertIn('our site', text)
        self.assertNotIn('](', text)

    def test_file_instruction_forbids_the_model_writing_its_own_link(self):
        instruction = _ai_document_instruction  # imported symbol still exists
        self.assertTrue(callable(instruction))
        from myapp.views import _ai_generated_file_instruction
        for name in ('generated.pdf', 'generated.docx', 'generated.txt'):
            with self.subTest(name=name):
                text = _ai_generated_file_instruction(name)
                self.assertIn('NEVER write a download link', text)
                self.assertIn('NO fenced block at all', text)


class ReportDrivenRoutingFixTests(TestCase):
    """Routing and prompt fixes traced to specific user reports."""

    def test_misspelt_edit_instruction_reaches_image_editing(self):
        # AIReport #42, verbatim. 'covert'/'postres' sent it to plain chat,
        # which answered "I cannot assist with that request".
        self.assertTrue(
            ai_chat.is_image_edit_instruction('COVERT INTO HIGH ANGAEMENT META ADS POSTRES')
        )

    def test_misspelt_generation_requests_still_route_to_image(self):
        for prompt in (
            'make a postre for my shop',
            'generate imge of a cat',
            'create a picutre of sunset',
            'walpaper of mountains banao',
        ):
            with self.subTest(prompt=prompt):
                self.assertTrue(ai_chat.is_image_generation_request(prompt))

    def test_analysis_requests_are_not_mistaken_for_edits(self):
        # These arrive with an image attached too. Treating them as edits would
        # send them to FLUX, which cannot answer a question about a picture.
        for prompt in (
            'what is in this image',
            'read the text in this photo',
            'explain this diagram',
            'summarise this document',
            'how much is the total in this bill',
        ):
            with self.subTest(prompt=prompt):
                self.assertFalse(ai_chat.is_image_edit_instruction(prompt))

    def test_ordinary_chat_is_not_pulled_into_image_generation(self):
        for prompt in (
            'hello how are you',
            'what is the capital of India',
            'write an email to my client',
            'explain recursion',
        ):
            with self.subTest(prompt=prompt):
                self.assertFalse(ai_chat.is_image_generation_request(prompt))

    def test_live_state_questions_trigger_a_web_search(self):
        # AIReport #43 answered "I don't have current operational status data"
        # with no search having run.
        for prompt in (
            'Status of Shree cement plant in meghalaya',
            'current status on the highway project',
            'is the factory still operational',
        ):
            with self.subTest(prompt=prompt):
                self.assertTrue(web_search.needs_search(prompt))

    def test_private_account_questions_never_cost_a_web_search(self):
        # A public search cannot answer these and would only add a round trip.
        for prompt in (
            'what is the status of my order',
            'my subscription status',
            'track my order',
            'make my whatsapp status funny',
        ):
            with self.subTest(prompt=prompt):
                self.assertFalse(web_search.needs_search(prompt))

    def test_prompt_states_the_capabilities_users_were_wrongly_denied(self):
        prompt = ai_chat.COMPACT_SYSTEM_PROMPT
        # #38: told users it could not display a generated image.
        self.assertIn('never say you cannot display', prompt)
        # #28/#33: claimed it could not create a file.
        self.assertIn('.docx', prompt)
        # #32: flat "I cannot help you with that" for a video request.
        self.assertIn('video/animation is not supported yet', prompt)
        # #42: refused a normal marketing request.
        self.assertIn('never refuse them', prompt)
        # #29: cited an EduTrellis page as the source of pharmacology facts.
        self.assertIn('never attach a', prompt)


class TruncatedOutputFollowUpTests(TestCase):
    """Reports #35 and #36: 'still it is half' truncated all over again.

    The follow-up carries no code or long-form keyword, so it fell back to the
    default token budget and cut off at the same place. The same user reported
    it twice.
    """

    def test_saying_the_answer_was_cut_short_earns_the_long_budget(self):
        for prompt in (
            'still it is half',
            'the code is not complete',
            'continue',
            'it got cut off',
            'yeh adhura hai',
            'aage likho',
            'baaki code do',
            'rest of the code please',
        ):
            with self.subTest(prompt=prompt):
                self.assertTrue(ai_chat.wants_long_form_output(prompt))
                self.assertTrue(ai_chat.is_truncated_output_complaint(prompt))

    def test_ordinary_turns_still_use_the_normal_budget(self):
        # The long budget also raises the request timeout, so this must not
        # fire on everyday chat.
        for prompt in (
            'hello',
            'what is the capital of India',
            'thanks',
            'write a tweet about coffee',
        ):
            with self.subTest(prompt=prompt):
                self.assertFalse(ai_chat.wants_long_form_output(prompt))

    def test_the_existing_long_form_cues_are_unaffected(self):
        # AIReport #29's original case must keep working.
        self.assertTrue(ai_chat.wants_long_form_output('Complete reference from Kdt'))
        self.assertFalse(ai_chat.is_truncated_output_complaint('give me a detailed breakdown'))


class DropboxArchiveIsOffDuringTestsTests(TestCase):
    """Regression guard: the suite must never upload into the live account.

    Several tests drive the real image and report views with only the storage
    layer mocked. Before DROPBOX_IMAGE_ARCHIVE_ENABLED existed, those quietly
    uploaded their fake images to the project owner's actual Dropbox.
    """

    def test_archiving_is_disabled_by_default_under_the_test_runner(self):
        self.assertFalse(dropbox_images.is_enabled())

    def test_a_test_run_cannot_queue_an_upload_with_the_real_credentials(self):
        with patch('myapp.dropbox_images._ensure_worker') as worker:
            self.assertFalse(dropbox_images.enqueue(b'bytes', 'png', 'real@example.com'))
            self.assertEqual(
                dropbox_images.enqueue_report_images(
                    1, 'real@example.com', user_image='data:image/png;base64,AA==',
                ),
                0,
            )

        worker.assert_not_called()

    def test_the_worker_itself_also_refuses_while_archiving_is_disabled(self):
        # The worker outlives any one request, so it re-checks rather than
        # trusting the decision made when the item was queued.
        with patch('myapp.dropbox_images.is_configured', return_value=True):
            self.assertIsNone(dropbox_images._build_client())


@override_settings(
    DROPBOX_IMAGE_ARCHIVE_ENABLED=True,
    DROPBOX_APP_KEY='test-key',
    DROPBOX_APP_SECRET='test-secret',
    DROPBOX_REFRESH_TOKEN='test-refresh-token',
)
class DropboxGeneratedImageArchiveTests(TestCase):
    """Generated images are mirrored to /vidhyora/<email>/ without delaying the reply."""

    PNG_BYTES = base64.b64decode(
        b'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=='
    )

    def setUp(self):
        cache.clear()
        # Module-level worker state survives between tests; start each one from
        # a drained queue and no cached client so assertions can't pick up a
        # previous test's upload.
        dropbox_images.flush(timeout=5)
        dropbox_images._client = None
        self.user = User.objects.create_user(
            username='archive-tests',
            email='Studio.Owner+AI@Example.com',
            password='test-password-123',
            is_staff=True,
        )
        self.client.force_login(self.user)

    def test_folder_name_is_lowercased_and_stripped_of_path_characters(self):
        # Dropbox rejects these characters outright in a path component, and
        # an unsanitised one would send the upload to a different folder.
        self.assertEqual(dropbox_images.folder_for('User@Example.com'), 'user@example.com')
        self.assertEqual(dropbox_images.folder_for('a/b\\c:d?e*f'), 'a_b_c_d_e_f')
        self.assertEqual(dropbox_images.folder_for('  spaced@x.com  '), 'spaced@x.com')

    def test_missing_email_falls_back_to_the_shared_guest_folder(self):
        # Guests can generate images too — they still get archived, just not
        # filed under an address that does not exist.
        for empty in ('', None, '   ', '...'):
            self.assertEqual(dropbox_images.folder_for(empty), dropbox_images.GUEST_FOLDER)

    def test_filename_is_sortable_and_keeps_the_local_storage_name(self):
        name = dropbox_images._filename('png', 'ai_generated/2026/09/05/abc123.png')

        self.assertTrue(name.endswith('-abc123.png'))
        # Leading YYYYMMDD-HHMMSS stamp, so a Dropbox folder listing sorts
        # chronologically by name.
        self.assertRegex(name, r'^\d{8}-\d{6}-')

    def test_unexpected_extension_is_not_trusted_into_the_filename(self):
        self.assertTrue(dropbox_images._filename('php', 'x.php').endswith('.png'))
        self.assertTrue(dropbox_images._filename('JPG', 'x.jpg').endswith('.jpg'))

    def test_enqueue_uploads_to_the_per_email_folder_under_vidhyora(self):
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            self.assertTrue(
                dropbox_images.enqueue(b'image-bytes', 'png', 'owner@example.com', 'local.png')
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        client.files_upload.assert_called_once()
        content, path = client.files_upload.call_args.args
        self.assertEqual(content, b'image-bytes')
        self.assertTrue(path.startswith('/vidhyora/owner@example.com/'))
        self.assertTrue(path.endswith('-local.png'))

    def test_upload_failures_are_swallowed_so_a_saved_image_is_never_lost(self):
        client = Mock()
        client.files_upload.side_effect = RuntimeError('Dropbox is down')
        with patch('myapp.dropbox_images._build_client', return_value=client):
            self.assertTrue(dropbox_images.enqueue(b'bytes', 'png', 'owner@example.com'))
            self.assertTrue(dropbox_images.flush(timeout=10))

        # A failed upload must also drop the cached client, so the next image
        # rebuilds one instead of reusing a possibly-broken connection.
        self.assertIsNone(dropbox_images._client)

    @override_settings(DROPBOX_REFRESH_TOKEN='')
    def test_incomplete_credentials_skip_the_upload_entirely(self):
        with patch('myapp.dropbox_images._build_client') as build:
            self.assertFalse(dropbox_images.enqueue(b'bytes', 'png', 'owner@example.com'))

        build.assert_not_called()

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/robot.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/2026/09/05/robot.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_generated_image_is_archived_under_the_logged_in_users_email(
        self, generate, save, storage_url,
    ):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')
        client = Mock()

        with patch('myapp.dropbox_images._build_client', return_value=client):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'A futuristic learning robot',
                    'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
                }),
                content_type='application/json',
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['X-Generated-Image-Url'], '/media/ai_generated/robot.png')

        content, path = client.files_upload.call_args.args
        self.assertEqual(content, self.PNG_BYTES)
        # The email is normalised to lowercase, so one person never ends up
        # with two archive folders.
        self.assertTrue(path.startswith('/vidhyora/studio.owner+ai@example.com/'))
        self.assertTrue(path.endswith('-robot.png'))

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/guest.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/guest.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_account_without_an_email_still_gets_its_image_archived(
        self, generate, save, storage_url,
    ):
        # An account can exist with a blank email (staff-created ones here do),
        # so the folder fallback has to cover logged-in users too, not just
        # guests — otherwise those images would go to '/vidhyora//...'.
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')
        no_email = User.objects.create_user(
            username='no-email-account', password='test-password-123', is_staff=True,
        )
        self.client.force_login(no_email)
        client = Mock()

        with patch('myapp.dropbox_images._build_client', return_value=client):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'A blue mountain',
                    'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
                }),
                content_type='application/json',
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertEqual(response.status_code, 200)
        _, path = client.files_upload.call_args.args
        self.assertTrue(path.startswith(f'/vidhyora/{dropbox_images.GUEST_FOLDER}/'))

    @patch('myapp.views.default_storage.url', return_value='/media/ai_generated/x.png')
    @patch('myapp.views.default_storage.save', return_value='ai_generated/x.png')
    @patch('myapp.views.image_generation.generate_image')
    def test_a_broken_dropbox_still_returns_the_image_to_the_user(
        self, generate, save, storage_url,
    ):
        generate.return_value = image_generation.GeneratedImage(self.PNG_BYTES, 'png')

        # The archive is a mirror, not the delivery path. Break it as badly as
        # possible — the real (unmocked) enqueue runs here, so this proves its
        # own error handling is what protects the reply, not the caller's.
        with patch('myapp.dropbox_images._ensure_worker', side_effect=RuntimeError('boom')):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'A red car',
                    'model': ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
                }),
                content_type='application/json',
            )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['X-Generated-Image-Url'], '/media/ai_generated/x.png')
        self.assertEqual(
            AIMessage.objects.get(role=AIMessage.ROLE_ASSISTANT).image_data,
            '/media/ai_generated/x.png',
        )

    def test_enqueue_never_raises_even_when_the_worker_cannot_start(self):
        with patch('myapp.dropbox_images._ensure_worker', side_effect=RuntimeError('no threads')):
            self.assertFalse(dropbox_images.enqueue(b'bytes', 'png', 'owner@example.com'))

    def _data_uri(self):
        return 'data:image/png;base64,' + base64.b64encode(self.PNG_BYTES).decode('ascii')

    def test_report_evidence_is_archived_under_the_reporters_reports_folder(self):
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            queued = dropbox_images.enqueue_report_images(
                42, 'Owner@Example.com',
                user_image=self._data_uri(), reply_image=self._data_uri(),
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertEqual(queued, 2)
        paths = sorted(call.args[1] for call in client.files_upload.call_args_list)
        self.assertEqual(len(paths), 2)
        for path in paths:
            self.assertTrue(path.startswith('/vidhyora/owner@example.com/reports/'))
        # Named by report number and side, so a reviewer can find the evidence
        # for report #42 without opening every file in the folder.
        self.assertTrue(paths[0].endswith('-report-42-reply.png'))
        self.assertTrue(paths[1].endswith('-report-42-user.png'))

    def test_report_evidence_data_uri_is_decoded_to_the_original_image_bytes(self):
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            dropbox_images.enqueue_report_images(7, 'owner@example.com', user_image=self._data_uri())
            self.assertTrue(dropbox_images.flush(timeout=10))

        # A real PNG must land in Dropbox, not the base64 text of one.
        content = client.files_upload.call_args.args[0]
        self.assertEqual(content, self.PNG_BYTES)
        self.assertTrue(content.startswith(b'\x89PNG'))

    def test_report_evidence_that_is_only_a_url_is_skipped_not_uploaded(self):
        # _snapshot_ai_report_image falls back to the bare media URL when the
        # file is already gone. There are no bytes behind that, so uploading it
        # would just create a file containing a URL.
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            queued = dropbox_images.enqueue_report_images(
                9, 'owner@example.com',
                user_image='/media/ai_generated/lost.png', reply_image='',
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertEqual(queued, 0)
        client.files_upload.assert_not_called()

    def test_corrupt_report_evidence_never_uploads_an_empty_file(self):
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            dropbox_images.enqueue_report_images(
                11, 'owner@example.com', user_image='data:image/png;base64,!!!not base64!!!',
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        client.files_upload.assert_not_called()

    def test_jpeg_report_evidence_keeps_a_usable_file_extension(self):
        client = Mock()
        with patch('myapp.dropbox_images._build_client', return_value=client):
            dropbox_images.enqueue_report_images(
                3, 'owner@example.com',
                user_image='data:image/jpeg;base64,' + base64.b64encode(b'\xff\xd8\xff-jpeg').decode(),
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertTrue(client.files_upload.call_args.args[1].endswith('.jpg'))

    def test_submitting_a_report_archives_its_image_evidence(self):
        conversation = AIConversation.objects.create(user=self.user, title='Report archive')
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_USER,
            content='Make this poster blue', image_data=self._data_uri(),
        )
        assistant = AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_ASSISTANT,
            content='', image_data=self._data_uri(),
            model_key=ai_chat.FLUX_KLEIN_4B_MODEL_KEY,
        )
        client = Mock()

        with patch('myapp.dropbox_images._build_client', return_value=client):
            response = self.client.post(
                '/AI/api/report/',
                data=json.dumps({
                    'conversation_id': conversation.id,
                    'message_id': assistant.pk,
                    'reply_image': assistant.image_data,
                    'explanation': 'The poster came out the wrong colour.',
                }),
                content_type='application/json',
            )
            self.assertTrue(dropbox_images.flush(timeout=10))

        self.assertEqual(response.status_code, 200)
        report = AIReport.objects.get()
        paths = sorted(call.args[1] for call in client.files_upload.call_args_list)
        self.assertEqual(len(paths), 2, f'expected both sides archived, got {paths}')
        for path in paths:
            self.assertTrue(
                path.startswith('/vidhyora/studio.owner+ai@example.com/reports/'), path,
            )
        self.assertTrue(paths[0].endswith(f'-report-{report.pk}-reply.png'))
        self.assertTrue(paths[1].endswith(f'-report-{report.pk}-user.png'))

    def test_a_broken_dropbox_still_lets_a_report_be_filed(self):
        conversation = AIConversation.objects.create(user=self.user, title='Report resilience')
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_USER, content='Why is this wrong?',
        )
        assistant = AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_ASSISTANT,
            content='A wrong answer.', model_key=ai_chat.CHATGPT_56_MODEL_KEY,
        )

        with patch('myapp.dropbox_images._ensure_worker', side_effect=RuntimeError('boom')):
            response = self.client.post(
                '/AI/api/report/',
                data=json.dumps({
                    'conversation_id': conversation.id,
                    'message_id': assistant.pk,
                    'reply_text': 'A wrong answer.',
                    'explanation': 'This is not correct.',
                }),
                content_type='application/json',
            )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(AIReport.objects.count(), 1)

    def test_image_archive_folder_cannot_collide_with_the_database_backup_folder(self):
        # dropbox_backup.py writes db.sqlite3 snapshots to its own root. If
        # these two ever shared a folder, delete_all_backups() would wipe the
        # users' generated images along with the backups.
        self.assertNotEqual(
            dropbox_images.ROOT_FOLDER.lower().rstrip('/'),
            dropbox_backup.BACKUP_ROOT.lower().rstrip('/'),
        )
        self.assertFalse(
            dropbox_backup.BACKUP_FOLDER.lower().startswith(
                dropbox_images.ROOT_FOLDER.lower().rstrip('/') + '/'
            )
        )


class AIResponseReliabilityTests(TestCase):
    def setUp(self):
        cache.clear()
        self.addCleanup(cache.clear)

    def test_request_routing_does_not_need_a_runtime_ml_model(self):
        self.assertEqual(request_router.classify('Debug this Python traceback'), 'code')
        self.assertEqual(request_router.classify('Research the latest facts and sources'), 'research')
        self.assertEqual(request_router.classify('Hello, how are you?'), 'general')
        self.assertEqual(request_router.choose_model('Fix this JavaScript bug', 'quick')[0], 'code')
        self.assertEqual(request_router.choose_chatgpt_worker('Hello, how are you?')[0], 'quick')
        self.assertEqual(request_router.choose_chatgpt_worker('Write a Python function')[0], 'code')

    def test_rate_limit_returns_friendly_user_message(self):
        user = User.objects.create_user(username='rate-limit@example.com', password='pw')
        self.client.force_login(user)

        with patch('myapp.views.AI_CHAT_RATE_LIMIT', 0):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({'message': 'hello'}),
                content_type='application/json',
            )

        self.assertEqual(response.status_code, 429)
        self.assertEqual(response.json()['status'], 'rate_limited')
        self.assertIn('Limit reached', response.json()['detail'])

    def test_chatgpt_uses_worker_model_with_stable_truthful_identity(self):
        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='answer'))])
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(
            create=Mock(return_value=iter([chunk])),
        )))

        with patch('myapp.ai_chat._get_client', return_value=client):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'Write a Python function'}],
                model_key='code', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))

        self.assertEqual(result, 'answer')
        request = client.chat.completions.create.call_args.kwargs
        self.assertEqual(request['model'], ai_chat.MODELS['code']['id'])
        self.assertEqual(request['messages'][0]['role'], 'system')
        self.assertIn('ChatGPT 5.6 in Vidhyora AI', request['messages'][0]['content'])
        self.assertIn('not the official OpenAI gpt-5.6 API', request['messages'][0]['content'])
        system_text = '\n'.join(
            item['content'] for item in request['messages'] if item['role'] == 'system'
        )
        self.assertIn('the only model name that may appear in your reply is ChatGPT 5.6', system_text)

    def test_chatgpt_identity_leak_is_caught_and_forced_to_a_safe_answer(self):
        """Live-observed: asked 'are you copy of gpt?' / 'who are you?', the
        ChatGPT 5.6 persona sometimes answered 'developed by researchers
        from NVIDIA' / 'trained by NVIDIA researchers' despite
        CHATGPT_56_SYSTEM_SUFFIX explicitly forbidding it. Every retry still
        leaking must end in the guaranteed-correct scripted answer, never
        the leaked text."""
        def make_stream(text):
            return iter([SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=text))])])

        # A fresh iterator per call — Mock(return_value=an_iterator) would
        # hand back the same already-exhausted iterator on every retry.
        create = Mock(side_effect=lambda **kw: make_stream('I was trained by NVIDIA researchers.'))
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client), patch('myapp.ai_chat.time.sleep'):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'are you copy of chatgpt?'}],
                model_key='quick', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))

        self.assertEqual(result, "I'm ChatGPT, developed by OpenAI.")
        self.assertNotIn('nvidia', result.lower())
        self.assertEqual(create.call_count, ai_chat.STREAM_RETRY_ATTEMPTS + 1)

    def test_chatgpt_identity_leak_caught_even_on_an_unrelated_question(self):
        """A real saved reply leaked 'trained by researchers from NVIDIA' as
        an unprompted aside answering 'do you knaow CodeXa Agency ??' — an
        explicit 'who are you' question never appeared. The guard checks
        every chatgpt56 reply's opening, not just ones that look like a
        provenance question, specifically to catch this."""
        def make_stream(text):
            return iter([SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=text))])])

        create = Mock(side_effect=lambda **kw: make_stream(
            "No, I'm a language model trained by researchers from NVIDIA."
        ))
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client), patch('myapp.ai_chat.time.sleep'):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'do you knaow CodeXa Agency ??'}],
                model_key='quick', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))

        self.assertEqual(result, "I'm ChatGPT, developed by OpenAI.")

    def test_chatgpt_identity_self_heals_when_a_later_retry_is_clean(self):
        def make_stream(text):
            return iter([SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=text))])])

        create = Mock(side_effect=[
            make_stream('Developed by researchers from NVIDIA.'),
            make_stream("I'm ChatGPT, developed by OpenAI."),
        ])
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client), patch('myapp.ai_chat.time.sleep'):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'who are you?'}],
                model_key='quick', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))

        self.assertEqual(result, "I'm ChatGPT, developed by OpenAI.")
        self.assertEqual(create.call_count, 2)

    def test_clean_chatgpt_reply_content_is_unchanged_short_or_long(self):
        """The opening-buffer check must never alter a clean reply's text —
        only its chunk boundaries (a short reply comes back as one combined
        chunk; a long one gets a buffered opening then streams the rest
        token-by-token same as before)."""
        short_words = ['Paris', ' is the capital.']
        create_short = Mock(return_value=iter([
            SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=c))])
            for c in short_words
        ]))
        client_short = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create_short)))
        with patch('myapp.ai_chat._get_client', return_value=client_short):
            result_short = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'what is the capital of france'}],
                model_key='quick', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))
        self.assertEqual(result_short, ''.join(short_words))

        long_words = ('Sure here is a detailed explanation of how TCP works ' * 15).split(' ')
        create_long = Mock(return_value=iter([
            SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=w + ' '))])
            for w in long_words
        ]))
        client_long = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create_long)))
        with patch('myapp.ai_chat._get_client', return_value=client_long):
            chunks_long = list(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'explain how tcp works'}],
                model_key='quick', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))
        self.assertEqual(''.join(chunks_long), ''.join(w + ' ' for w in long_words))
        # Opening buffered as one block, then streamed per-token afterward —
        # not one giant chunk, and not unbuffered from the very first token.
        self.assertGreater(len(chunks_long), 1)
        self.assertGreaterEqual(len(chunks_long[0]), ai_chat.IDENTITY_CHECK_BUFFER_CHARS)

    def test_chatgpt_with_attached_image_checks_both_vision_and_identity_without_duplicating_text(self):
        """A chatgpt56 turn with an attached image gets routed to the vision
        worker internally while staying identified as ChatGPT 5.6, so both
        check_vision_opening and check_identity_opening are active on the
        same reply — they used to share one buffer variable, which meant
        the vision-buffered opening got yielded once normally and then
        appended into the identity buffer and yielded a second time."""
        parts = ['I can see ', 'a cat ', 'in this image. ' * 20, 'END']
        create = Mock(return_value=iter([
            SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content=c))]) for c in parts
        ]))
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': [
                    {'type': 'image_url', 'image_url': {'url': 'data:image/png;base64,xx'}},
                    {'type': 'text', 'text': 'what is this'},
                ]}],
                model_key='vision', identity_model_key=ai_chat.CHATGPT_56_MODEL_KEY,
            ))

        self.assertEqual(result, ''.join(parts))

    def test_transient_failure_on_non_default_model_falls_back_to_quick(self):
        """A busy Vision/Ultra/Code worker should still get a real answer via
        Quick instead of surfacing a hard failure — AIReport #10 and #30 both
        showed a non-default model just failing outright with no fallback at
        all, because the fallback used to be gated on model_key ==
        DEFAULT_MODEL_KEY (which is 'chatgpt56', not the model actually used
        here), making it dead for every other model."""
        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='answer'))])
        create = Mock(side_effect=[TimeoutError('request timed out'), iter([chunk])])
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client), \
             patch('myapp.ai_chat.time.sleep'):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'describe this image'}], model_key='vision',
            ))

        self.assertEqual(result, 'answer')
        self.assertEqual(create.call_count, 2)
        first_model = create.call_args_list[0].kwargs['model']
        second_model = create.call_args_list[1].kwargs['model']
        self.assertEqual(first_model, ai_chat.MODELS['vision']['id'])
        self.assertEqual(second_model, ai_chat.MODELS['quick']['id'])

    def test_long_form_request_gets_larger_token_budget_and_timeout(self):
        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='answer'))])
        create = Mock(return_value=iter([chunk]))
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))

        with patch('myapp.ai_chat._get_client', return_value=client):
            ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'Complete reference on cholinergic drugs'}],
                model_key='quick', max_tokens=6000,
            ))

        request = create.call_args.kwargs
        self.assertEqual(request['max_tokens'], 6000)
        self.assertEqual(request['timeout'], ai_chat.STREAM_TIMEOUT_LONG)

    def test_image_request_detection_against_real_reported_failures(self):
        """Locks in fixes for real AIReport prompts that used to be
        misrouted — see the AI Reports dashboard analysis. Each of these
        used to reach a text model instead of the image pipeline."""
        # AIReport #7, #8, #9, #17 — verb+noun, including Hinglish word order.
        for prompt in [
            'Change image background to white background and its sise in 1000x1000px',
            'Hey char GPT can use generate images',
            'Cat ka image bna ke do',
            'create design of Nashik360 logo',
        ]:
            self.assertTrue(ai_chat.is_image_generation_request(prompt), prompt)

        # AIReport #31, #34 — informal noun/verb abbreviations.
        self.assertTrue(ai_chat.is_image_generation_request('plz gen img'))
        self.assertTrue(ai_chat.is_image_generation_request('make motor drawing'))

        # AIReport #19 — bare descriptive prompt, no generate verb at all.
        self.assertTrue(ai_chat.is_probable_image_prompt(
            'A realistic brown dog sitting on a grassy field, shiny coat, '
            'bright eyes, soft lighting, high detail, 4k.'
        ))

        # AIReport #11 — edit instruction on an attached image with no
        # image-shaped noun ("Names" isn't one).
        self.assertTrue(ai_chat.is_image_edit_instruction('REmove all Names'))

        # AIReport #29 — explicit long-form ask that used to be cut off.
        self.assertTrue(ai_chat.wants_long_form_output('Complete reference from Kdt'))

        # AIReport #27 — "birthday card" wasn't recognised; only the fixed
        # two-word phrase "greeting card" was in the noun list.
        self.assertTrue(ai_chat.is_image_generation_request('please make a birthday card'))

        # AIReport #45 — an attached-image "turn in to instagram post" (note
        # the two-word "in to", not "into") fell through to Vision because
        # neither "post" nor a "turn into" verb form was recognised.
        self.assertTrue(ai_chat.is_image_generation_request('turn in to instagram post'))
        self.assertTrue(ai_chat.is_image_edit_instruction('turn in to instagram post'))

        # AIReport #46 — "Made" (past tense of "make") wasn't in the verb list.
        self.assertTrue(ai_chat.is_image_generation_request('Made light background of this post'))

        # AIReport #39 — "use this logo" on an attached image is a real
        # composite/edit instruction with no verb from the edit-only list.
        self.assertTrue(ai_chat.is_image_edit_instruction('use this logo'))
        # But a generic "use this ..." with no image-shaped noun must not
        # misfire on an attached image that's actually about something else.
        self.assertFalse(ai_chat.is_image_edit_instruction('use this data to build a report'))

        # AIReport #44 — Hindi "Is ladki ko cafe me dikhao" ("show/place this
        # girl in a cafe") on an attached photo needs an edit-shaped verb;
        # none of dikhao/daalo/lagao/jodo/nikaal were recognised before.
        self.assertTrue(ai_chat.is_image_edit_instruction('Is ladki ko cafe me dikhao'))

        # Must not misfire on ordinary chat/analysis text.
        for prompt in [
            'what is in this image', 'describe this photo', 'is this a cat or a dog',
            'write a short story about a dog in a field, it should be heartwarming',
        ]:
            self.assertFalse(ai_chat.is_image_edit_instruction(prompt), prompt)
            self.assertFalse(ai_chat.is_probable_image_prompt(prompt), prompt)

    def test_chatgpt_is_the_fresh_default_on_every_page_load(self):
        staff = User.objects.create_user(
            username='default-model-staff@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(staff)
        response = self.client.get('/AI/')

        self.assertEqual(ai_chat.DEFAULT_MODEL_KEY, ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response.context['ai_default_model'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response.context['ai_default_model_label'], 'ChatGPT 5.6')
        self.assertNotContains(response, "localStorage.getItem('ai_model')")
        self.assertNotContains(response, "localStorage.setItem('ai_model'")

    def test_free_users_only_get_quick_and_code(self):
        user = User.objects.create_user(username='free-models@example.com', password='test-password-123')
        StoreProfile.objects.create(user=user)
        self.client.force_login(user)

        page = self.client.get('/AI/')
        self.assertEqual(page.context['ai_default_model'], 'quick')
        access = {item['key']: item['locked'] for item in page.context['ai_models']}
        self.assertFalse(access['quick'])
        self.assertNotIn('light', access)
        self.assertFalse(access['code'])
        self.assertTrue(access[ai_chat.CHATGPT_56_MODEL_KEY])
        self.assertTrue(access['ultra'])
        self.assertTrue(access['reasoning'])
        self.assertTrue(access[ai_chat.FLUX_KLEIN_4B_MODEL_KEY])
        self.assertContains(page, 'Free users can use Quick and Code.')

        blocked = self.client.post(
            '/AI/api/send/',
            data=json.dumps({'message': 'Use the premium model', 'model': 'ultra'}),
            content_type='application/json',
        )
        self.assertEqual(blocked.status_code, 403)
        self.assertEqual(blocked.json()['status'], 'subscription_required')
        self.assertEqual(AIConversation.objects.filter(user=user).count(), 0)

    def test_free_quick_is_allowed_but_locked_automatic_image_routing_is_not(self):
        user = User.objects.create_user(username='free-quick@example.com', password='test-password-123')
        StoreProfile.objects.create(user=user)
        self.client.force_login(user)

        with patch('myapp.views.ai_chat.stream_chat', return_value=iter(['Quick reply'])):
            allowed = self.client.post(
                '/AI/api/send/',
                data=json.dumps({'message': 'Hello', 'model': 'quick'}),
                content_type='application/json',
            )
            self.assertEqual(allowed.status_code, 200)
            self.assertEqual(b''.join(allowed.streaming_content).decode(), 'Quick reply')

        blocked_image = self.client.post(
            '/AI/api/send/',
            data=json.dumps({'message': 'Generate an image of a mountain', 'model': 'quick'}),
            content_type='application/json',
        )
        self.assertEqual(blocked_image.status_code, 403)
        self.assertEqual(blocked_image.json()['status'], 'subscription_required')

    def test_premium_user_keeps_all_models(self):
        user = User.objects.create_user(username='premium-models@example.com', password='test-password-123')
        StoreProfile.objects.create(
            user=user, ai_subscription_until=timezone.now() + timedelta(days=30),
        )
        self.client.force_login(user)

        page = self.client.get('/AI/')
        self.assertEqual(page.context['ai_default_model'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertFalse(any(item['locked'] for item in page.context['ai_models']))

        with patch('myapp.views.ai_chat.stream_chat', return_value=iter(['Premium reply'])):
            allowed = self.client.post(
                '/AI/api/send/',
                data=json.dumps({'message': 'Solve this carefully', 'model': 'ultra'}),
                content_type='application/json',
            )
            self.assertEqual(allowed.status_code, 200)
            self.assertEqual(b''.join(allowed.streaming_content).decode(), 'Premium reply')

    def test_chatgpt_routes_general_code_and_image_turns(self):
        user = User.objects.create_user(
            username='chatgpt-router@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)

        with patch('myapp.views.ai_chat.stream_chat', side_effect=lambda *args, **kwargs: iter(['reply'])) as stream_chat:
            with patch('myapp.views.image_ocr.extract_data_uri', return_value=''):
                cases = (
                    ({'message': 'Hello there'}, 'quick', 'general'),
                    ({'message': 'Debug this Python function'}, 'code', 'code'),
                    ({'message': 'What is in this?', 'image': 'data:image/png;base64,AA=='}, 'vision', 'image'),
                )
                for extra_payload, worker_key, category in cases:
                    payload = {'model': ai_chat.CHATGPT_56_MODEL_KEY, **extra_payload}
                    response = self.client.post(
                        '/AI/api/send/', data=json.dumps(payload), content_type='application/json',
                    )
                    self.assertEqual(response.status_code, 200)
                    self.assertEqual(b''.join(response.streaming_content).decode(), 'reply')
                    self.assertEqual(response['X-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
                    self.assertEqual(response['X-Routed-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
                    self.assertEqual(response['X-Request-Category'], category)
                    call = stream_chat.call_args
                    self.assertEqual(call.kwargs['model_key'], worker_key)
                    self.assertEqual(call.kwargs['identity_model_key'], ai_chat.CHATGPT_56_MODEL_KEY)

    def test_chatgpt_reply_cannot_expose_a_worker_name_split_across_chunks(self):
        user = User.objects.create_user(
            username='chatgpt-identity-lock@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)

        leaked_chunks = iter([
            'I am generating images using FL',
            'UX.2 Klein 4B model through NVIDIA Nemotron. ',
            'Vidhyora Code helped too.',
        ])
        with patch('myapp.views.ai_chat.stream_chat', return_value=leaked_chunks):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'Can you generate images?',
                    'model': ai_chat.CHATGPT_56_MODEL_KEY,
                }),
                content_type='application/json',
            )
            body = b''.join(response.streaming_content).decode()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response['X-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertEqual(response['X-Routed-Model-Key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertIn('ChatGPT 5.6', body)
        for hidden_name in ('FLUX', 'NVIDIA', 'Nemotron', 'Vidhyora Code'):
            self.assertNotIn(hidden_name.lower(), body.lower())
        assistant = AIMessage.objects.get(role=AIMessage.ROLE_ASSISTANT)
        self.assertEqual(assistant.content, body)
        self.assertEqual(assistant.model_key, ai_chat.CHATGPT_56_MODEL_KEY)

    def test_chatgpt_never_claims_a_backend_vendor_trained_it(self):
        """The reported symptom: replies saying "I was trained by NVIDIA".

        ai_chat's own retry guard only inspects the opening few hundred
        characters, so a claim made partway through a long answer used to
        reach the browser untouched.
        """
        from myapp.views import _chatgpt_public_reply

        for leak in (
            'I was trained by NVIDIA.',
            'My underlying model was developed by NVIDIA.',
            "I'm an NVIDIA model.",
            'I was trained by Meta on the Llama architecture.',
            'I am based on the Nemotron base model from NVIDIA.',
            'I was developed by NVIDIA, not OpenAI.',
            'A' * 600 + ' To be clear, I was actually built by NVIDIA.',
        ):
            cleaned = _chatgpt_public_reply(leak)
            for vendor in ('nvidia', 'nemotron', 'llama', 'mistral'):
                self.assertNotIn(vendor, cleaned.lower(), msg=leak[:60])

        # ...while a genuine answer *about* those companies must survive: the
        # word itself is not the problem, claiming it built this assistant is.
        for factual in (
            'NVIDIA is a semiconductor company founded in 1993.',
            'GPUs are made by NVIDIA and AMD.',
            'Llama is an open-weights model family released by Meta.',
        ):
            self.assertEqual(_chatgpt_public_reply(factual), factual)

    def test_chatgpt_streams_progressively_without_duplicating_text(self):
        """ChatGPT 5.6 used to withhold the whole reply until generation
        finished — nothing rendered for the entire wait. It now releases text
        as it arrives, holding back only a short tail for the sanitizer."""
        user = User.objects.create_user(
            username='chatgpt-streaming@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)

        sentence = 'This is a normal, clean sentence of assistant output. '
        chunks = [sentence] * 40
        with patch('myapp.views.ai_chat.stream_chat', return_value=iter(chunks)):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'Explain something at length',
                    'model': ai_chat.CHATGPT_56_MODEL_KEY,
                }),
                content_type='application/json',
            )
            parts = [part.decode() for part in response.streaming_content]

        body = ''.join(parts)
        expected = sentence * 40
        # Delivered exactly once, in full, and in order.
        self.assertEqual(body, expected)
        # Genuinely progressive: the reply arrived in several pieces rather
        # than one final dump, and the first piece came well before the end.
        released = [p for p in parts if p]
        self.assertGreater(len(released), 1)
        self.assertLess(len(released[0]), len(expected))
        assistant = AIMessage.objects.get(role=AIMessage.ROLE_ASSISTANT)
        self.assertEqual(assistant.content, expected)

    def test_chatgpt_reports_disconnected_text_access_without_worker_names(self):
        class RemovedWorkerError(Exception):
            status_code = 404

        user = User.objects.create_user(
            username='chatgpt-disconnected@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        upstream = RemovedWorkerError(
            "NVIDIA Function FLUX/Nemotron worker not found for account",
        )
        with patch('myapp.views.ai_chat.stream_chat', side_effect=upstream):
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({'message': 'hi', 'model': ai_chat.CHATGPT_56_MODEL_KEY}),
                content_type='application/json',
            )
            body = b''.join(response.streaming_content).decode()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(
            body,
            'ChatGPT 5.6 text access is currently disconnected. '
            'Update the configured text-model API key, then restart the application.',
        )
        for hidden_name in ('NVIDIA', 'FLUX', 'Nemotron'):
            self.assertNotIn(hidden_name.lower(), body.lower())
        self.assertFalse(AIMessage.objects.filter(role=AIMessage.ROLE_ASSISTANT).exists())

    def test_explicit_file_request_is_routed_and_downloadable_from_every_model(self):
        cache.clear()
        self.addCleanup(cache.clear)
        user = User.objects.create_user(
            username='file-owner@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        prompt = 'Generate a file named greeting.txt with the exact content Hello world and share a download link.'

        with patch('myapp.views.ai_chat.stream_chat', side_effect=lambda *args, **kwargs: iter(['```txt\nHello world\n```'])) as stream_chat:
            for selected_model in ai_chat.MODELS:
                with self.subTest(model=selected_model):
                    response = self.client.post(
                        '/AI/api/send/',
                        data=json.dumps({'message': prompt, 'model': selected_model}),
                        content_type='application/json',
                    )
                    self.assertEqual(response.status_code, 200)
                    body = b''.join(response.streaming_content).decode()
                    self.assertIn('[Download greeting.txt](', body)
                    expected_public_route = (
                        ai_chat.CHATGPT_56_MODEL_KEY
                        if selected_model == ai_chat.CHATGPT_56_MODEL_KEY
                        else 'code'
                    )
                    self.assertEqual(response['X-Routed-Model-Key'], expected_public_route)
                    self.assertEqual(response['X-Request-Category'], 'file_generation')
                    self.assertIn('application, not you', stream_chat.call_args.kwargs['document_instruction'])

        self.assertEqual(AIGeneratedFile.objects.filter(user=user).count(), len(ai_chat.MODELS))
        generated_file = AIGeneratedFile.objects.filter(user=user).first()
        self.assertEqual(generated_file.file_name, 'greeting.txt')
        self.assertEqual(generated_file.content, 'Hello world')

        download = self.client.get(f'/AI/api/files/{generated_file.token}/download/')
        self.assertEqual(download.status_code, 200)
        self.assertEqual(download.content.decode(), 'Hello world')
        self.assertEqual(download['Content-Disposition'], 'attachment; filename="greeting.txt"')
        self.assertEqual(download['Cache-Control'], 'private, no-store')

        other_user = User.objects.create_user(username='other-file-user@example.com', password='pw')
        self.client.force_login(other_user)
        self.assertEqual(self.client.get(f'/AI/api/files/{generated_file.token}/download/').status_code, 404)

    def test_generated_file_intent_and_fence_extraction_are_conservative(self):
        self.assertEqual(
            _ai_generated_file_spec('Create a Python file named hello.py with a print statement.'),
            {'file_name': 'hello.py'},
        )
        self.assertEqual(
            _ai_generated_file_spec('Prepare a downloadable markdown document.'),
            {'file_name': 'generated.md'},
        )
        self.assertEqual(
            _ai_generated_file_spec('make a summarise note in word file'),
            {'file_name': 'generated.docx'},
        )
        self.assertIsNone(_ai_generated_file_spec('Explain what this Python file does.'))
        self.assertEqual(
            _extract_ai_generated_file_content('```python\nprint("hello")\n```'),
            'print("hello")',
        )

    def test_report_28_returns_a_genuine_downloadable_word_document(self):
        cache.clear()
        self.addCleanup(cache.clear)
        user = User.objects.create_user(
            username='word-file-owner@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)

        with patch(
            'myapp.views.ai_chat.stream_chat',
            return_value=iter(['```markdown\n# Summary Note\n\n- First important point\n- Second important point\n```']),
        ) as stream_chat:
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': 'make a summarise note in word file',
                    'model': ai_chat.CHATGPT_56_MODEL_KEY,
                }),
                content_type='application/json',
            )
            body = b''.join(response.streaming_content).decode()

        self.assertEqual(response.status_code, 200)
        self.assertIn('[Download generated.docx](', body)
        self.assertIn('genuine DOCX file', stream_chat.call_args.kwargs['document_instruction'])
        generated_file = AIGeneratedFile.objects.get(user=user)
        download = self.client.get(f'/AI/api/files/{generated_file.token}/download/')

        self.assertEqual(download.status_code, 200)
        self.assertEqual(
            download['Content-Type'],
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        )
        self.assertTrue(download.content.startswith(b'PK'))
        word_document = doc_extract.DocxDocument(io.BytesIO(download.content))
        document_text = '\n'.join(paragraph.text for paragraph in word_document.paragraphs)
        self.assertIn('Summary Note', document_text)
        self.assertIn('First important point', document_text)

    def test_report_33_pdf_request_with_attached_image_still_generates_a_real_pdf(self):
        # AIReport #33: "Make renewal notice to send to client using this
        # data and add logo I have attached in pdf" — an attached image
        # used to unconditionally skip file-generation routing and fall
        # through to Vision (which can only describe an image, never
        # produce a download), and PDF wasn't even a supported output
        # format yet. This locks in both fixes: the object regex accepts
        # a bare "pdf" (not just "file"/"document"), file-generation intent
        # is checked even with an image attached, and the download is a
        # genuine PDF. Embedding the attached logo into the PDF itself is
        # not implemented — only the real text content and download link.
        cache.clear()
        self.addCleanup(cache.clear)
        user = User.objects.create_user(
            username='pdf-file-owner@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        prompt = 'Make renewal notice to send to client using this data and add logo I have attached in pdf'

        with patch(
            'myapp.views.ai_chat.stream_chat',
            return_value=iter(['```markdown\n# Renewal Notice\n\n- Policy is due for renewal\n```']),
        ) as stream_chat:
            response = self.client.post(
                '/AI/api/send/',
                data=json.dumps({
                    'message': prompt,
                    'image': 'data:image/png;base64,AA==',
                    'model': ai_chat.CHATGPT_56_MODEL_KEY,
                }),
                content_type='application/json',
            )
            body = b''.join(response.streaming_content).decode()

        self.assertEqual(response.status_code, 200)
        self.assertIn('[Download generated.pdf](', body)
        self.assertEqual(response['X-Request-Category'], 'file_generation')
        self.assertIn('genuine PDF', stream_chat.call_args.kwargs['document_instruction'])
        generated_file = AIGeneratedFile.objects.get(user=user)
        self.assertEqual(generated_file.file_name, 'generated.pdf')

        download = self.client.get(f'/AI/api/files/{generated_file.token}/download/')
        self.assertEqual(download.status_code, 200)
        self.assertEqual(download['Content-Type'], 'application/pdf')
        self.assertTrue(download.content.startswith(b'%PDF'))
        pdf_text = ''.join(page.extract_text() for page in PdfReader(io.BytesIO(download.content)).pages)
        self.assertIn('Renewal Notice', pdf_text)
        self.assertIn('Policy is due for renewal', pdf_text)

    def test_every_model_is_told_the_real_current_date_and_time(self):
        """A model can't read a clock, so "what's today's date?" was answered
        from its training cutoff. The live clock is now stated on every turn."""
        note = ai_chat.current_datetime_note()
        now = datetime.datetime.now(ZoneInfo('Asia/Kolkata'))
        self.assertIn(now.strftime('%d %B %Y'), note)
        self.assertIn(str(now.year), note)
        self.assertIn('IST', note)
        self.assertIn("never say you don't have access to the current date", note.lower())

        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='ok'))])
        for model_key in ('quick', 'ultra', 'code', 'vision', 'reasoning', ai_chat.CHATGPT_56_MODEL_KEY):
            create = Mock(return_value=iter([chunk]))
            client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(create=create)))
            with patch('myapp.ai_chat._get_client', return_value=client):
                list(ai_chat.stream_chat(
                    [{'role': 'user', 'content': 'what is the date today?'}],
                    model_key=model_key,
                ))
            system_prompt = create.call_args.kwargs['messages'][0]['content']
            self.assertIn(now.strftime('%d %B %Y'), system_prompt, msg=model_key)

    def test_web_search_only_fires_on_time_sensitive_questions(self):
        for should_search in (
            'what is the latest news about AI',
            'current gold rate in india',
            'who won the match yesterday',
            'aaj ka petrol price kya hai',
            'search for the best hosting providers',
        ):
            self.assertTrue(web_search.needs_search(should_search), msg=should_search)

        # A search is a network round trip on the critical path of a reply, so
        # everything answerable without one must stay out of it.
        for should_not in (
            'write a python function to sort a list',
            'rephrase this message for me',
            'generate an image of a cat',
            'what is 15% of 2400',
            'explain object oriented programming',
            'hello how are you',
        ):
            self.assertFalse(web_search.needs_search(should_not), msg=should_not)

    def test_web_search_failure_degrades_to_a_normal_answer(self):
        """A search outage must never break the chat — it just means the model
        answers from its own knowledge, as it did before search existed."""
        cache.clear()
        self.addCleanup(cache.clear)
        # ddgs.DDGS is a lazy proxy that forwards to ddgs.ddgs.DDGS; patching
        # the proxy leaves the real class (and so the real network call) in
        # place, so the implementation class is the target here.
        with patch('ddgs.ddgs.DDGS.text', side_effect=RuntimeError('rate limited')):
            self.assertEqual(web_search.search('current gold rate'), [])
            self.assertIsNone(web_search.build_context('current gold rate'))

    def test_web_results_are_passed_to_the_model_as_grounding(self):
        cache.clear()
        self.addCleanup(cache.clear)
        fake = [{'title': 'Gold Rate Today', 'href': 'https://example.com/gold',
                 'body': 'Gold is 71,000 per 10g today.'}]
        with patch('ddgs.ddgs.DDGS.text', return_value=fake):
            context = web_search.build_context('current gold rate in india')
        self.assertIn('Gold Rate Today', context)
        self.assertIn('https://example.com/gold', context)
        self.assertIn('71,000', context)
        # The model must be told not to invent beyond what was actually found.
        self.assertIn('never invent a result', context.lower())

    def test_spreadsheet_and_presentation_requests_produce_real_office_files(self):
        for message, expected in (
            ('make an excel sheet of monthly expenses', 'generated.xlsx'),
            ('give me a slide deck about our services', 'generated.pptx'),
            ('put this data in a spreadsheet', 'generated.xlsx'),
            ('mujhe ek presentation chahiye', 'generated.pptx'),
        ):
            self.assertEqual(_ai_generated_file_spec(message), {'file_name': expected}, msg=message)

        workbook_bytes = _ai_excel_bytes('Item,Qty,Price\nPens,10,25.5\n"Books, hardcover",3,499\n')
        self.assertTrue(workbook_bytes.startswith(b'PK'))
        sheet = load_workbook(io.BytesIO(workbook_bytes)).active
        rows = list(sheet.iter_rows(values_only=True))
        self.assertEqual(rows[0], ('Item', 'Qty', 'Price'))
        # Numbers stored as numbers, so the sheet is actually usable for
        # formulas, and a quoted value containing a comma stays one cell.
        self.assertEqual(rows[1], ('Pens', 10, 25.5))
        self.assertEqual(rows[2][0], 'Books, hardcover')

        deck_bytes = _ai_powerpoint_bytes('# Intro\n- Who we are\n# Services\n- SEO\n- Websites\n')
        self.assertTrue(deck_bytes.startswith(b'PK'))
        deck = Presentation(io.BytesIO(deck_bytes))
        self.assertEqual([slide.shapes.title.text for slide in deck.slides], ['Intro', 'Services'])
        second = [p.text for p in deck.slides[1].placeholders[1].text_frame.paragraphs]
        self.assertEqual(second, ['SEO', 'Websites'])

    def test_file_conversion_covers_every_offered_format_pair(self):
        """Every pair the UI offers must actually produce a valid file — an
        offered conversion that then fails is worse than not offering it."""
        docx_source = file_convert.text_to_docx_bytes('# Report\n\n- one\n- two\n\nA paragraph.')
        pdf_source = file_convert.text_to_pdf_bytes('# Invoice\n\n- line A\n\nTotal 4999.')
        csv_source = b'Item,Qty,Price\nPens,10,25.5\n"Books, hardcover",3,499\n'
        xlsx_source = file_convert.rows_to_xlsx_bytes([['Item', 'Qty'], ['Pens', '10']])
        image_buffer = io.BytesIO()
        Image.new('RGBA', (120, 80), (255, 0, 0, 128)).save(image_buffer, 'PNG')
        png_source = image_buffer.getvalue()

        signatures = {
            'pdf': b'%PDF', 'docx': b'PK', 'xlsx': b'PK',
            'jpg': b'\xff\xd8\xff', 'png': b'\x89PNG', 'webp': b'RIFF',
        }
        sources = {
            'invoice.pdf': pdf_source, 'report.docx': docx_source,
            'data.csv': csv_source, 'sheet.xlsx': xlsx_source,
            'logo.png': png_source, 'notes.txt': b'plain text\nsecond line',
        }
        for name, data in sources.items():
            targets = file_convert.targets_for(name)
            self.assertTrue(targets, msg=name)
            for target in targets:
                payload, filename, _ = file_convert.convert(data, name, target)
                self.assertTrue(payload, msg=f'{name}->{target}')
                self.assertTrue(filename.endswith(f'.{target}'), msg=f'{name}->{target}')
                if target in signatures:
                    self.assertTrue(
                        payload.startswith(signatures[target]), msg=f'{name}->{target}',
                    )

        # Content actually survives the round trip, rather than producing a
        # valid-but-empty file.
        csv_out, _, _ = file_convert.convert(xlsx_source, 'sheet.xlsx', 'csv')
        self.assertIn(b'Item', csv_out)
        xlsx_out, _, _ = file_convert.convert(csv_source, 'data.csv', 'xlsx')
        rows = list(load_workbook(io.BytesIO(xlsx_out)).active.iter_rows(values_only=True))
        self.assertEqual(rows[1], ('Pens', 10, 25.5))
        self.assertEqual(rows[2][0], 'Books, hardcover')

    def test_file_conversion_rejects_unsupported_pairs_with_a_clear_reason(self):
        with self.assertRaises(file_convert.ConvertError):
            file_convert.convert(b'data', 'thing.exe', 'pdf')
        with self.assertRaises(file_convert.ConvertError):
            file_convert.convert(b'data', 'notes.txt', 'xlsx')
        with self.assertRaises(file_convert.ConvertError):
            file_convert.convert(b'', 'notes.txt', 'pdf')
        # A scanned PDF has no text layer; say so instead of returning an
        # empty document that looks like a successful conversion.
        blank_pdf = file_convert.text_to_pdf_bytes('')
        with self.assertRaises(file_convert.ConvertError) as caught:
            file_convert.convert(blank_pdf, 'scan.pdf', 'docx')
        self.assertIn('scan', str(caught.exception).lower())

    def test_convert_endpoint_returns_the_converted_file(self):
        cache.clear()
        self.addCleanup(cache.clear)
        user = User.objects.create_user(
            username='convert-user@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        upload = SimpleUploadedFile('data.csv', b'Name,Score\nAsha,91\n', content_type='text/csv')

        response = self.client.post('/AI/api/convert/', {'file': upload, 'target': 'xlsx'})
        self.assertEqual(response.status_code, 200)
        self.assertTrue(response.content.startswith(b'PK'))
        self.assertIn('data.xlsx', response['Content-Disposition'])
        self.assertEqual(response['Cache-Control'], 'private, no-store')

        bad = SimpleUploadedFile('data.csv', b'Name,Score\n', content_type='text/csv')
        rejected = self.client.post('/AI/api/convert/', {'file': bad, 'target': 'docx'})
        self.assertEqual(rejected.status_code, 400)
        self.assertIn('CSV', rejected.json()['detail'])

    def test_conversation_exports_as_a_real_pdf_and_word_file(self):
        user = User.objects.create_user(
            username='export-user@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        conversation = AIConversation.objects.create(user=user, title='Pricing questions')
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_USER, content='What are your rates?',
        )
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_ASSISTANT,
            content='Our website packages start at 14999.',
        )

        pdf = self.client.get(f'/AI/api/conversations/{conversation.id}/export/pdf/')
        self.assertEqual(pdf.status_code, 200)
        self.assertEqual(pdf['Content-Type'], 'application/pdf')
        self.assertTrue(pdf.content.startswith(b'%PDF'))
        text = ''.join(page.extract_text() for page in PdfReader(io.BytesIO(pdf.content)).pages)
        self.assertIn('What are your rates?', text)
        self.assertIn('14999', text)

        docx = self.client.get(f'/AI/api/conversations/{conversation.id}/export/docx/')
        self.assertEqual(docx.status_code, 200)
        self.assertTrue(docx.content.startswith(b'PK'))

        self.assertEqual(
            self.client.get(f'/AI/api/conversations/{conversation.id}/export/rtf/').status_code, 400,
        )
        # Someone else's conversation must not be exportable.
        other = User.objects.create_user(username='other-export@example.com', password='pw')
        self.client.force_login(other)
        self.assertEqual(
            self.client.get(f'/AI/api/conversations/{conversation.id}/export/pdf/').status_code, 404,
        )

    def test_report_on_an_image_only_turn_still_records_what_was_asked(self):
        """Six real reports arrived with a blank prompt because the reported
        turn carried only an image, leaving nothing to diagnose."""
        user = User.objects.create_user(
            username='report-context@example.com', password='test-password-123', is_staff=True,
        )
        self.client.force_login(user)
        conversation = AIConversation.objects.create(user=user, title='Poster help')
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_USER,
            content='make me a birthday poster',
        )
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_USER,
            content='', image_data='data:image/png;base64,AA==',
        )
        reply = AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_ASSISTANT,
            content='I cannot help with that.', model_key=ai_chat.CHATGPT_56_MODEL_KEY,
        )

        response = self.client.post(
            '/AI/api/report/',
            data=json.dumps({
                'conversation_id': conversation.id, 'message_id': reply.id,
                'reply_text': 'I cannot help with that.',
                'explanation': 'poster nahin ban raha',
            }),
            content_type='application/json',
        )
        self.assertEqual(response.status_code, 200)
        report = AIReport.objects.get()
        # The blank turn is described, and the real instruction is carried
        # through so the report can actually be triaged.
        self.assertIn('image with no text', report.user_prompt)
        self.assertIn('make me a birthday poster', report.user_prompt)

    def test_accuracy_rules_cover_maths_and_unclear_images(self):
        self.assertTrue(ai_chat.is_math_request('Solve 2x + 5 = 17'))
        self.assertTrue(ai_chat.is_math_request('Calculate 18% of 450'))
        self.assertFalse(ai_chat.is_math_request('Write a friendly customer email'))
        self.assertTrue(ai_chat.is_code_request('Fix this Django traceback'))
        self.assertFalse(ai_chat.is_code_request('Write a friendly customer email'))
        self.assertIn('never give only a number', ai_chat.COMPACT_SYSTEM_PROMPT)
        self.assertIn('ask for a clearer image', ai_chat.COMPACT_SYSTEM_PROMPT)
        self.assertIn('Never claim an action', ai_chat.COMPACT_SYSTEM_PROMPT)
        self.assertIn('complete, secure, directly usable code', ai_chat.CODE_SYSTEM_SUFFIX)
        self.assertIn('Do not claim code was executed or tested', ai_chat.CODE_SYSTEM_SUFFIX)

        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='answer'))])
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(
            create=Mock(return_value=iter([chunk])),
        )))
        with patch('myapp.ai_chat._get_client', return_value=client):
            list(ai_chat.stream_chat([{'role': 'user', 'content': 'Calculate 2 + 3'}], model_key='quick'))
        sent_messages = client.chat.completions.create.call_args.kwargs['messages']
        late_reminder = sent_messages[-2]['content']
        self.assertIn('step-by-step', late_reminder)
        self.assertIn('**Final answer:**', late_reminder)
        self.assertIn('verify the result', late_reminder)

    def test_mixed_multimodal_request_gets_complete_response_rules(self):
        prompt = (
            "1. Calculate 15% of 800.\n"
            "2. Analyze the attached screenshot.\n"
            "3. Fix this Python error.\n"
            "4. Explain the relevant Django setting.\n"
            "5. Check whether the logic is valid.\n"
            "6. Product cost is 500, advertising is 100, selling price is 900; calculate profit percentage."
        )
        self.assertEqual(ai_chat.count_user_requests(prompt), 6)
        chunk = SimpleNamespace(choices=[SimpleNamespace(delta=SimpleNamespace(content='answer'))])
        client = SimpleNamespace(chat=SimpleNamespace(completions=SimpleNamespace(
            create=Mock(return_value=iter([chunk])),
        )))
        content = [
            {'type': 'text', 'text': prompt},
            {'type': 'image_url', 'image_url': {'url': 'data:image/png;base64,AA=='}},
        ]
        with patch('myapp.ai_chat._get_client', return_value=client):
            list(ai_chat.stream_chat([{'role': 'user', 'content': content}], model_key='vision'))

        sent_messages = client.chat.completions.create.call_args.kwargs['messages']
        reminder = sent_messages[-2]['content']
        self.assertIn('Answer every one exactly once', reminder)
        self.assertIn('numbered section per item', reminder)
        self.assertIn('Total cost = product cost + advertising expense', reminder)
        self.assertIn('net profit / total cost × 100', reminder)
        self.assertIn('frontend-safe Markdown', reminder)
        self.assertIn('never labels like pythonCopy', reminder)
        self.assertIn('Analyse the attached image itself', reminder)
        self.assertIn('continue answering all other items', reminder)

    def test_note_router_understands_numbered_read_and_edit_commands(self):
        self.assertEqual(request_router.match_read_note('open note 1'), '1')
        self.assertEqual(request_router.match_read_note('read my note #2'), '#2')
        self.assertEqual(request_router.match_edit_note('edit note 1'), ('1', ''))
        self.assertEqual(request_router.match_edit_note('edit note 1 to Call at 7'), ('1', 'Call at 7'))
        self.assertTrue(request_router.is_note_intent('create a note: Call at 7'))

    def test_note_router_tolerates_common_typos_without_over_correcting(self):
        # A missed match here doesn't fail quietly — it falls through to the
        # real AI model, which (per its own history of this router's past
        # confirmations) fabricates its own fake "done!" instead of just not
        # understanding. See request_router._typo_correct_note_keywords.
        self.assertEqual(request_router.match_delete_note('delet all notyes'), request_router.DELETE_ALL_NOTES)
        self.assertTrue(request_router.is_note_intent('tkae this noet'))
        self.assertTrue(request_router.is_show_notes_intent('shwo my notess'))
        self.assertEqual(request_router.match_edit_note('edti note about milk to bread'), ('milk', 'bread'))
        self.assertEqual(request_router.match_read_note('opne note 1'), '1')
        self.assertTrue(request_router.is_note_intent('remmember a noet: call home'))
        self.assertEqual(request_router.match_delete_note('eraze note 2'), '2')
        self.assertEqual(request_router.match_edit_note('renmae note 1 to New title'), ('1', 'New title'))
        self.assertTrue(request_router.is_show_notes_intent('reed all notyes'))
        # Ordinary sentences that merely contain a word close to one of the
        # trigger keywords must never get swept in as a false positive —
        # 'made' -> 'make' would otherwise turn a past-tense remark into a
        # live "make a note" command.
        self.assertFalse(request_router.is_note_intent('she made a note about it yesterday, what should I do'))
        self.assertFalse(request_router.is_note_intent('I have not opened the store today'))

    def test_company_context_contains_verified_ai_contacts(self):
        self.assertTrue(company_knowledge.is_company_query('what is the sales team number?'))
        self.assertIn('+91 96959 53183', company_knowledge.PUBLIC_SITE_CONTEXT)
        self.assertIn('Vidhyora AI is an AI assistant', company_knowledge.PUBLIC_SITE_CONTEXT)
        self.assertNotIn('/websitecreation', company_knowledge.PUBLIC_SITE_CONTEXT)
        self.assertNotIn('/store', company_knowledge.PUBLIC_SITE_CONTEXT)

    def test_company_query_detection_covers_realistic_contact_phrasings(self):
        # These specific phrasings are what actually reached the AI model
        # ungrounded before this fix (the old regex required 'your ...' or
        # 'company's ...' or the brand name) — that gap, not a wrong fact in
        # the prompt itself, is what let it fabricate a fake US toll-free
        # number and a fake sales@edutrellis.com email for "sales team
        # number" style questions. See business_info.py for the real values.
        for query in (
            'sales number', 'contact number', 'WhatsApp number', 'sales email',
            'contact EduTrellis', 'company address', 'customer-support email',
            'how can I contact you?',
        ):
            self.assertTrue(company_knowledge.is_company_query(query), msg=query)
        # Unrelated messages must not get swept in as a false positive.
        for query in ('explain object oriented programming', 'help me write a poem'):
            self.assertFalse(company_knowledge.is_company_query(query), msg=query)

    def test_no_fabricated_contact_details_anywhere_in_ai_facing_text(self):
        wrong_markers = ('555', 'edutrellis.com', 'sales@edutrellis', '1-800', '1‑800')
        for text in (ai_chat.SYSTEM_PROMPT, company_knowledge.PUBLIC_SITE_CONTEXT):
            for marker in wrong_markers:
                self.assertNotIn(marker, text)
        # The real values must come from one shared source, not be retyped.
        self.assertIn(business_info.PHONE_DISPLAY, ai_chat.SYSTEM_PROMPT)
        self.assertIn(business_info.EMAIL_SUPPORT, ai_chat.SYSTEM_PROMPT)
        self.assertIn(business_info.PHONE_DISPLAY, company_knowledge.PUBLIC_SITE_CONTEXT)
        self.assertIn('no separate sales line', ai_chat.SYSTEM_PROMPT)
        self.assertIn('toll-free', company_knowledge.PUBLIC_SITE_CONTEXT)

    @override_settings(AI_USE_PRESIDIO=False)
    def test_fast_privacy_path_redacts_common_identifiers(self):
        redacted = privacy.redact('Email me@example.com or call 9876543210')

        self.assertEqual(redacted, 'Email <EMAIL> or call <PHONE>')

    def test_default_model_retries_on_quick_backend(self):
        chunk = SimpleNamespace(
            choices=[SimpleNamespace(delta=SimpleNamespace(content='Recovered reply'))]
        )
        create = SimpleNamespace()
        create.create = Mock(
            side_effect=[TimeoutError('upstream timed out'), iter([chunk])]
        )
        client = SimpleNamespace(chat=SimpleNamespace(completions=create))

        with patch('myapp.ai_chat._get_client', return_value=client):
            result = ''.join(ai_chat.stream_chat(
                [{'role': 'user', 'content': 'Hello'}], model_key=ai_chat.DEFAULT_MODEL_KEY
            ))

        self.assertEqual(result, 'Recovered reply')
        self.assertEqual(create.create.call_count, 2)
        self.assertEqual(
            create.create.call_args_list[1].kwargs['model'],
            ai_chat.MODELS['quick']['id'],
        )


class AINoteCRUDTests(TestCase):
    """Exercises My Notes end-to-end through /AI/api/send/, not just the
    request_router regex parsing — request_router.match_read_note/
    match_edit_note being correct in isolation once shipped with a plain
    `re.fullmatch(...)` call in views._ai_matching_notes with no `import re`
    at the top of views.py, which 500'd every read/edit/delete-by-number
    request while create/show (which never call that function) kept working
    silently. A regex-level unit test alone can't catch that class of bug."""
    def setUp(self):
        self.user = User.objects.create_user(
            username='note-crud@example.com', email='note-crud@example.com', password='test-password-123',
        )
        StoreProfile.objects.create(user=self.user, phone='9999999999')
        self.client.force_login(self.user)

    def send(self, message, conversation_id=None):
        payload = {'message': message}
        if conversation_id:
            payload['conversation_id'] = conversation_id
        response = self.client.post('/AI/api/send/', data=json.dumps(payload), content_type='application/json')
        body = b''.join(response.streaming_content).decode('utf-8')
        self.assertEqual(response.status_code, 200, msg=body)
        return response, body

    def test_full_note_lifecycle_via_chat(self):
        response, _ = self.send('note down: buy milk and eggs')
        conversation_id = int(response['X-Conversation-Id'])
        self.send('note down: call dentist tomorrow', conversation_id)
        self.assertEqual(AINote.objects.filter(user=self.user).count(), 2)

        _, show_body = self.send('show my notes', conversation_id)
        self.assertIn('buy milk and eggs', show_body)
        self.assertIn('call dentist tomorrow', show_body)

        _, read_body = self.send('open note 1', conversation_id)
        self.assertIn('call dentist tomorrow', read_body)

        edit_response, edit_body = self.send('replace note about milk with buy milk, eggs and bread', conversation_id)
        self.assertEqual(edit_response['X-Notes-Changed'], '1')
        self.assertIn('buy milk, eggs and bread', AINote.objects.get(heading__icontains='bread').content)

        delete_response, delete_body = self.send('delete note about dentist', conversation_id)
        self.assertEqual(delete_response['X-Notes-Changed'], '1')
        self.assertEqual(AINote.objects.filter(user=self.user).count(), 1)

        self.send('delete all my notes', conversation_id)
        self.assertEqual(AINote.objects.filter(user=self.user).count(), 0)

    def test_contextual_partial_edit_preserves_the_rest_of_the_note(self):
        response, _ = self.send('add note: I have to work tomorow at 4pm')
        conversation_id = int(response['X-Conversation-Id'])
        note = AINote.objects.get(user=self.user)
        self.assertEqual(note.content, 'I have to work tomorrow at 4pm')

        self.send('edit note 1', conversation_id)
        update_response, update_body = self.send('update the time to 7pm', conversation_id)
        self.assertEqual(update_response['X-Notes-Changed'], '1')
        note.refresh_from_db()
        self.assertEqual(note.content, 'I have to work tomorrow at 7pm')
        self.assertIn(note.content, update_body)

    def test_ambiguous_contextual_edit_asks_then_applies_the_choice(self):
        response, _ = self.send('save note: Meeting tomorrow at 4pm')
        conversation_id = int(response['X-Conversation-Id'])
        self.send('update the first note', conversation_id)

        ambiguous_response, ambiguous_body = self.send('update it to 7pm', conversation_id)
        self.assertNotIn('X-Notes-Changed', ambiguous_response)
        self.assertEqual(ambiguous_body, 'Should I update only the time to 7pm, or replace the full note?')
        self.assertEqual(AINote.objects.get(user=self.user).content, 'Meeting tomorrow at 4pm')

        final_response, final_body = self.send('only the time', conversation_id)
        self.assertEqual(final_response['X-Notes-Changed'], '1')
        self.assertEqual(AINote.objects.get(user=self.user).content, 'Meeting tomorrow at 7pm')
        self.assertIn('Meeting tomorrow at 7pm', final_body)

    def test_explicit_database_id_targets_note_without_text_search(self):
        response, _ = self.send('add note: Original text')
        conversation_id = int(response['X-Conversation-Id'])
        note = AINote.objects.get(user=self.user)

        update_response, _ = self.send(f'replace note id {note.pk} with Replaced by database ID', conversation_id)
        self.assertEqual(update_response['X-Notes-Changed'], '1')
        note.refresh_from_db()
        self.assertEqual(note.content, 'Replaced by database ID')

    def test_rename_changes_only_heading_and_never_stores_sidebar_number(self):
        response, _ = self.send('write note: Client meeting at 3pm')
        conversation_id = int(response['X-Conversation-Id'])
        note = AINote.objects.get(user=self.user)

        rename_response, rename_body = self.send('rename the first note to Tomorrow meeting', conversation_id)
        self.assertEqual(rename_response['X-Notes-Changed'], '1')
        note.refresh_from_db()
        self.assertEqual(note.heading, 'Tomorrow meeting')
        self.assertEqual(note.content, 'Client meeting at 3pm')
        self.assertFalse(note.heading.startswith('1.'))
        self.assertIn('Tomorrow meeting', rename_body)
        self.assertIn('Client meeting at 3pm', rename_body)

    def test_bare_save_never_copies_chat_history_and_empty_copy_is_exact(self):
        response, body = self.send('take a note')
        self.assertEqual(AINote.objects.filter(user=self.user).count(), 0)
        self.assertEqual(body, 'What would you like the note to say?')

        conversation_id = int(response['X-Conversation-Id'])
        _, body = self.send('show all notes', conversation_id)
        self.assertEqual(body, 'You don’t have any saved notes.')

    def test_chat_listing_and_sidebar_api_use_identical_database_snapshot(self):
        response, _ = self.send('create a note: first database note')
        conversation_id = int(response['X-Conversation-Id'])
        self.send('remember a note: second database note', conversation_id)

        api_notes = self.client.get('/AI/api/notes/').json()['notes']
        _, chat_body = self.send('view all notes', conversation_id)
        self.assertEqual([item['content'] for item in api_notes], ['second database note', 'first database note'])
        for item in api_notes:
            self.assertIn(item['content'], chat_body)
            self.assertEqual(chat_body.count(item['content']), 1)

        delete_response, _ = self.send('delet all notyes', conversation_id)
        self.assertEqual(delete_response['X-Notes-Changed'], '1')
        self.assertEqual(self.client.get('/AI/api/notes/').json()['notes'], [])


class AIAccountProfileTests(TestCase):
    def setUp(self):
        self.media_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.media_dir.cleanup)
        media_override = override_settings(MEDIA_ROOT=self.media_dir.name)
        media_override.enable()
        self.addCleanup(media_override.disable)

        self.user = User.objects.create_user(
            username='account@example.com', email='account@example.com',
            password='old-password', first_name='Account', last_name='Owner',
        )
        self.profile = StoreProfile.objects.create(
            user=self.user, phone='9999999999',
            phone_verified=True,
            ai_subscription_until=timezone.now() + timedelta(days=10),
        )
        self.other_user = User.objects.create_user(
            username='other@example.com', email='other@example.com', password='test-password-123',
        )
        StoreProfile.objects.create(user=self.other_user, phone='8888888888')
        self.client.force_login(self.user)

    def test_account_api_returns_subscription_and_only_own_report_statuses(self):
        conversation = AIConversation.objects.create(user=self.user, title='My reported chat')
        AIReport.objects.create(
            user=self.user, conversation=conversation, reported_reply='Incorrect reply',
            explanation='The answer was wrong.', model_key='quick', status=AIReport.STATUS_RESOLVED,
        )
        other_conversation = AIConversation.objects.create(user=self.other_user, title='Private other chat')
        AIReport.objects.create(
            user=self.other_user, conversation=other_conversation, reported_reply='Other reply',
            explanation='Must remain private.', status=AIReport.STATUS_OPEN,
        )

        response = self.client.get('/AI/api/account/')
        body = response.json()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(body['user']['name'], 'Account Owner')
        self.assertEqual(body['subscription']['plan_name'], 'Vidhyora AI Premium')
        self.assertTrue(body['subscription']['active'])
        self.assertEqual(len(body['reports']), 1)
        self.assertEqual(body['reports'][0]['status'], 'resolved')
        self.assertEqual(body['reports'][0]['status_label'], 'Resolved')
        self.assertNotContains(response, 'Private other chat')
        self.assertEqual(response['Cache-Control'], 'private, no-store')

    def test_report_submit_snapshots_the_preceding_user_question(self):
        conversation = AIConversation.objects.create(user=self.user, title='Chat')
        AIMessage.objects.create(conversation=conversation, role=AIMessage.ROLE_USER, content='What is 2+2?')
        AIMessage.objects.create(
            conversation=conversation, role=AIMessage.ROLE_ASSISTANT,
            content='It is 5.', model_key='quick',
        )

        response = self.client.post(
            '/AI/api/report/', content_type='application/json',
            data=json.dumps({
                'conversation_id': conversation.id,
                'reply_text': 'It is 5.',
                'model_key': 'quick',
                'explanation': 'Wrong answer.',
            }),
        )

        self.assertEqual(response.status_code, 200)
        report = AIReport.objects.get(conversation=conversation)
        self.assertEqual(report.user_prompt, 'What is 2+2?')
        self.assertEqual(report.reported_reply, 'It is 5.')

    def test_profile_update_changes_login_email_name_phone_and_avatar(self):
        image_bytes = io.BytesIO()
        Image.new('RGB', (20, 20), (220, 20, 45)).save(image_bytes, format='PNG')
        avatar = SimpleUploadedFile('avatar.png', image_bytes.getvalue(), content_type='image/png')

        response = self.client.post('/AI/api/profile/update/', {
            'name': 'Updated Person', 'email': 'updated@example.com',
            'phone': '7777777777', 'avatar': avatar,
        })
        self.user.refresh_from_db()
        self.profile.refresh_from_db()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(self.user.get_full_name(), 'Updated Person')
        self.assertEqual(self.user.email, 'updated@example.com')
        self.assertEqual(self.user.username, 'account@example.com')
        self.assertEqual(self.profile.phone, '7777777777')
        self.assertFalse(self.profile.phone_verified)
        self.assertTrue(self.profile.avatar.name.endswith('.png'))
        self.assertContains(response, 'avatar_url')

    def test_profile_update_rejects_another_accounts_email(self):
        response = self.client.post('/AI/api/profile/update/', {
            'name': 'Account Owner', 'email': 'other@example.com', 'phone': '9999999999',
        })

        self.assertEqual(response.status_code, 400)
        self.assertIn('email', response.json()['errors'])
        self.user.refresh_from_db()
        self.assertEqual(self.user.email, 'account@example.com')

    def test_password_change_keeps_user_logged_in(self):
        response = self.client.post(
            '/AI/api/profile/password/',
            data=json.dumps({'current_password': 'old-password', 'new_password': 'new-password'}),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.user.refresh_from_db()
        self.assertTrue(self.user.check_password('new-password'))
        self.assertEqual(self.client.get('/AI/api/account/').status_code, 200)

    def test_ai_page_contains_profile_dropdown_and_all_account_panels(self):
        response = self.client.get('/AI/')

        self.assertContains(response, 'Edit profile')
        self.assertContains(response, 'Upload profile image')
        self.assertContains(response, 'Subscription details')
        self.assertContains(response, 'My reports')
        self.assertNotContains(response, '> Admin Panel</a>')

    def test_superuser_profile_dropdown_contains_admin_panel(self):
        superuser = User.objects.create_superuser(
            username='superadmin@example.com', email='superadmin@example.com', password='admin-password',
        )
        self.client.force_login(superuser)

        response = self.client.get('/AI/')

        self.assertContains(response, '> Admin Panel</a>')
        self.assertContains(response, 'href="/store/dashboard/"')
        self.assertNotContains(response, 'Back to site')

    def test_anonymous_user_cannot_read_account_details(self):
        self.client.logout()

        response = self.client.get('/AI/api/account/')

        self.assertEqual(response.status_code, 401)


class GitHubAccessTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username='github-user@example.com', email='github-user@example.com',
            password='github-password', is_staff=False,
        )
        StoreProfile.objects.create(user=self.user, phone='9555555555')

    def test_github_button_and_status_are_available_to_regular_logged_in_user(self):
        self.client.force_login(self.user)

        page = self.client.get('/AI/')
        status = self.client.get('/AI/api/github/status/')

        self.assertEqual(page.status_code, 200)
        self.assertContains(page, 'id="githubBtn"')
        self.assertContains(page, 'Connect a repository')
        self.assertEqual(status.status_code, 200)
        self.assertEqual(status.json(), {'status': 'ok', 'connected': False})

        self.client.logout()
        guest_page = self.client.get('/AI/')
        guest_status = self.client.get('/AI/api/github/status/')
        self.assertNotContains(guest_page, 'id="githubBtn"')
        self.assertEqual(guest_status.status_code, 401)

    @patch('myapp.views.github_ops.list_user_repos')
    @patch('myapp.views.github_ops.get_authenticated_user')
    def test_regular_user_can_connect_token_and_select_repository(self, get_user, list_repos):
        get_user.return_value = {'login': 'octocat'}
        list_repos.return_value = [
            {'full_name': 'octocat/demo', 'private': False, 'default_branch': 'main'},
            {'full_name': 'octocat/private-app', 'private': True, 'default_branch': 'develop'},
        ]
        self.client.force_login(self.user)

        response = self.client.post(
            '/AI/api/github/connect/', data=json.dumps({'token': 'secret-test-token'}),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.assertNotContains(response, 'secret-test-token')
        connection = GitHubConnection.objects.get(user=self.user)
        self.assertEqual(connection.github_username, 'octocat')
        self.assertEqual(connection.repo_full_name, 'octocat/demo')
        self.assertEqual(connection.access_token, 'secret-test-token')

        with patch('myapp.views.github_ops.get_repo', return_value={
            'full_name': 'octocat/private-app', 'default_branch': 'develop',
        }):
            select = self.client.post(
                '/AI/api/github/repo/', data=json.dumps({'repo': 'octocat/private-app'}),
                content_type='application/json',
            )
        self.assertEqual(select.status_code, 200)
        connection.refresh_from_db()
        self.assertEqual(connection.repo_full_name, 'octocat/private-app')
        self.assertEqual(connection.default_branch, 'develop')

    @override_settings(GITHUB_OAUTH_CLIENT_ID='github-client-id')
    def test_regular_user_can_start_github_oauth(self):
        self.client.force_login(self.user)

        response = self.client.get('/AI/api/github/oauth/start/')

        self.assertEqual(response.status_code, 302)
        self.assertTrue(response.url.startswith('https://github.com/login/oauth/authorize?'))
        self.assertIn('github_oauth_state', self.client.session)

    @override_settings(
        GITHUB_OAUTH_CLIENT_ID='github-client-id',
        GITHUB_OAUTH_CLIENT_SECRET='github-client-secret',
    )
    @patch('myapp.views.github_ops.list_user_repos')
    @patch('myapp.views.github_ops.get_authenticated_user')
    @patch('myapp.views.requests.post')
    def test_regular_user_can_complete_github_oauth(self, post, get_user, list_repos):
        token_response = Mock()
        token_response.json.return_value = {'access_token': 'oauth-secret-token'}
        post.return_value = token_response
        get_user.return_value = {'login': 'oauth-user'}
        list_repos.return_value = [
            {'full_name': 'oauth-user/project', 'private': True, 'default_branch': 'main'},
        ]
        self.client.force_login(self.user)
        session = self.client.session
        session['github_oauth_state'] = 'expected-state'
        session.save()

        response = self.client.get(
            '/AI/api/github/oauth/callback/?state=expected-state&code=temporary-code',
        )

        self.assertRedirects(response, '/AI/?github_connected=1', fetch_redirect_response=False)
        connection = GitHubConnection.objects.get(user=self.user)
        self.assertEqual(connection.access_token, 'oauth-secret-token')
        self.assertEqual(connection.github_username, 'oauth-user')
        self.assertEqual(connection.repo_full_name, 'oauth-user/project')

    @patch('myapp.views.ai_chat.github_plan_changes')
    @patch('myapp.views.ai_chat.github_select_files')
    @patch('myapp.views.github_ops.create_pull_request')
    @patch('myapp.views.github_ops.upsert_file')
    @patch('myapp.views.github_ops.create_branch')
    @patch('myapp.views.github_ops.get_branch_sha')
    @patch('myapp.views.github_ops.get_file')
    @patch('myapp.views.github_ops.get_tree')
    def test_regular_user_prompt_pushes_review_branch_and_opens_pull_request(
        self, get_tree, get_file, get_branch_sha, create_branch, upsert_file,
        create_pull_request, select_files, plan_changes,
    ):
        GitHubConnection.objects.create(
            user=self.user, access_token='secret-token', github_username='octocat',
            repo_full_name='octocat/demo', default_branch='main',
        )
        get_tree.return_value = ['app.py', 'README.md']
        select_files.return_value = ['app.py']
        get_file.return_value = ('print("old")\n', 'existing-sha')
        plan_changes.return_value = {
            'summary': 'Updated the greeting.',
            'commit_message': 'Update greeting',
            'operations': [
                {'action': 'update', 'path': 'app.py', 'content': 'print("hello")\n'},
            ],
        }
        get_branch_sha.return_value = 'base-sha'
        create_pull_request.return_value = {'html_url': 'https://github.com/octocat/demo/pull/7'}
        self.client.force_login(self.user)

        response = self.client.post('/AI/api/github/send/', data=json.dumps({
            'message': 'Change the greeting in app.py',
            'model': ai_chat.CHATGPT_56_MODEL_KEY,
        }), content_type='application/json')

        self.assertEqual(response.status_code, 200)
        body = response.json()
        self.assertEqual(body['status'], 'ok')
        self.assertEqual(body['model_key'], ai_chat.CHATGPT_56_MODEL_KEY)
        self.assertIn('https://github.com/octocat/demo/pull/7', body['reply'])
        created_branch = create_branch.call_args.args[3]
        self.assertTrue(created_branch.startswith('ai/'))
        upsert_file.assert_called_once_with(
            'secret-token', 'octocat', 'demo', 'app.py', 'print("hello")\n',
            'Update greeting', created_branch, sha='existing-sha',
        )
        create_pull_request.assert_called_once()
        self.assertEqual(AIMessage.objects.filter(conversation__user=self.user).count(), 2)
        self.assertEqual(
            AIMessage.objects.get(
                conversation__user=self.user, role=AIMessage.ROLE_ASSISTANT,
            ).model_key,
            ai_chat.CHATGPT_56_MODEL_KEY,
        )


class LocationConsentTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username='location-user@example.com', email='location-user@example.com',
            password='location-password',
        )
        self.profile = StoreProfile.objects.create(user=self.user, phone='9444444444')

    def test_authenticated_user_can_save_one_location_fix(self):
        self.client.force_login(self.user)

        response = self.client.post('/AI/api/location/', data=json.dumps({
            'consent': 'granted', 'latitude': 20.296059,
            'longitude': 85.824539, 'accuracy': 18.6,
        }), content_type='application/json')

        self.assertEqual(response.status_code, 200)
        self.profile.refresh_from_db()
        self.assertEqual(self.profile.location_consent, StoreProfile.LOCATION_GRANTED)
        self.assertEqual(self.profile.location_latitude, Decimal('20.296059'))
        self.assertEqual(self.profile.location_longitude, Decimal('85.824539'))
        self.assertEqual(self.profile.location_accuracy_m, 19)
        self.assertIsNotNone(self.profile.location_updated_at)

    def test_declining_location_is_saved_and_clears_old_coordinates(self):
        self.profile.location_consent = StoreProfile.LOCATION_GRANTED
        self.profile.location_latitude = Decimal('20.296059')
        self.profile.location_longitude = Decimal('85.824539')
        self.profile.location_accuracy_m = 20
        self.profile.save()
        self.client.force_login(self.user)

        response = self.client.post(
            '/AI/api/location/', data=json.dumps({'consent': 'denied'}),
            content_type='application/json',
        )

        self.assertEqual(response.status_code, 200)
        self.profile.refresh_from_db()
        self.assertEqual(self.profile.location_consent, StoreProfile.LOCATION_DENIED)
        self.assertIsNone(self.profile.location_latitude)
        self.assertIsNone(self.profile.location_longitude)
        self.assertIsNone(self.profile.location_accuracy_m)

    def test_invalid_or_anonymous_location_updates_are_rejected(self):
        anonymous = self.client.post(
            '/AI/api/location/', data=json.dumps({'consent': 'denied'}),
            content_type='application/json',
        )
        self.assertEqual(anonymous.status_code, 401)

        self.client.force_login(self.user)
        invalid = self.client.post('/AI/api/location/', data=json.dumps({
            'consent': 'granted', 'latitude': 95, 'longitude': 85, 'accuracy': 10,
        }), content_type='application/json')
        self.assertEqual(invalid.status_code, 400)
        self.profile.refresh_from_db()
        self.assertEqual(self.profile.location_consent, StoreProfile.LOCATION_UNKNOWN)

    def test_prompt_is_shown_until_a_choice_has_been_saved(self):
        self.client.force_login(self.user)

        response = self.client.get('/')
        self.assertTrue(response.context['show_location_prompt'])
        self.assertContains(response, 'Enable one-time access')
        self.assertContains(response, 'No continuous or background tracking')

        self.profile.location_consent = StoreProfile.LOCATION_DENIED
        self.profile.location_updated_at = timezone.now()
        self.profile.save(update_fields=['location_consent', 'location_updated_at'])
        response = self.client.get('/')
        self.assertFalse(response.context['show_location_prompt'])


class AIDashboardOverviewTests(TestCase):
    def setUp(self):
        self.admin = User.objects.create_superuser(
            username='dashboard-admin@example.com', email='dashboard-admin@example.com',
            password='admin-password',
        )
        StoreProfile.objects.create(user=self.admin, phone='9000000000')
        self.customer = User.objects.create_user(
            username='ai-customer@example.com', email='ai-customer@example.com', password='customer-password',
        )
        self.customer_profile = StoreProfile.objects.create(
            user=self.customer, phone='9111111111',
            manual_amount_paid=Decimal('1250.50'),
            ai_subscription_until=timezone.now() + timedelta(days=30),
            location_consent=StoreProfile.LOCATION_GRANTED,
            location_latitude=Decimal('20.296059'),
            location_longitude=Decimal('85.824539'),
            location_accuracy_m=25,
            location_updated_at=timezone.now(),
        )
        self.customer_chat = AIConversation.objects.create(user=self.customer, title='Customer AI chat')
        self.guest_chat = AIConversation.objects.create(
            session_key='guest-session', ip_address='203.0.113.10', title='Guest AI chat',
        )
        AIMessage.objects.create(conversation=self.customer_chat, role=AIMessage.ROLE_USER, content='Question')
        AIMessage.objects.create(
            conversation=self.customer_chat, role=AIMessage.ROLE_ASSISTANT,
            content='Answer', model_key='quick',
        )
        AIReport.objects.create(
            user=self.customer, conversation=self.customer_chat, reported_reply='Answer',
            explanation='Needs correction', status=AIReport.STATUS_OPEN,
        )
        AIReport.objects.create(
            conversation=self.guest_chat, session_key='guest-session', reported_reply='Guest answer',
            explanation='Already handled', status=AIReport.STATUS_RESOLVED,
        )
        AIBlock.objects.create(user=self.customer, reason='Test block', created_by=self.admin)
        self.client.force_login(self.admin)

    def test_overview_contains_ai_metrics_and_recent_ai_data(self):
        response = self.client.get('/store/dashboard/')

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.context['total_conversations'], 2)
        self.assertEqual(response.context['total_messages'], 2)
        self.assertEqual(response.context['registered_ai_users'], 1)
        self.assertEqual(response.context['guest_conversations'], 1)
        self.assertEqual(response.context['active_subscribers'], 1)
        self.assertEqual(response.context['open_reports'], 1)
        self.assertEqual(response.context['resolved_reports'], 1)
        self.assertEqual(response.context['active_blocks'], 1)
        self.assertContains(response, 'Customer AI chat')
        self.assertContains(response, 'Guest AI chat')
        self.assertContains(response, 'Needs correction')
        self.assertContains(response, 'AI Overview')
        self.assertEqual(len(response.context['daily_activity']), 7)
        self.assertEqual(response.context['registered_conversations'], 1)
        self.assertEqual(response.context['registered_pct'], 50)
        self.assertEqual(response.context['guest_pct'], 50)
        self.assertEqual(response.context['report_total'], 2)
        self.assertEqual(response.context['report_open_pct'], 50)
        self.assertEqual(response.context['model_usage'][0]['key'], 'quick')
        self.assertEqual(response.context['location_enabled'], 1)
        self.assertEqual(response.context['location_declined'], 0)
        self.assertEqual(response.context['location_not_asked'], 1)
        self.assertEqual(response.context['location_enabled_pct'], 50)
        self.assertContains(response, 'AI activity')
        self.assertContains(response, 'Conversation audience')
        self.assertContains(response, 'Report status')
        self.assertContains(response, 'AI model usage')
        self.assertContains(response, 'Location permission')

    def test_sidebar_keeps_only_ai_pwa_and_backup_options(self):
        response = self.client.get('/store/dashboard/')

        for label in ('Overview', 'Signups', 'AI Management', 'AI Activity', 'AI Reports', 'PWA / Install App', 'Backup &amp; Restore'):
            self.assertContains(response, label)
        self.assertContains(response, 'href="/" class="dash-logo"')
        self.assertNotContains(response, 'Back to store')
        self.assertNotContains(response, 'Back to site')
        for removed_path in (
            '/store/dashboard/contacts/', '/store/dashboard/categories/',
            '/store/dashboard/products/', '/store/dashboard/orders/', '/store/dashboard/delivery/',
            '/store/dashboard/payments/', '/store/dashboard/payment-settings/',
            '/store/dashboard/fee-settings/', '/store/dashboard/email-settings/',
            '/store/dashboard/about/', '/store/dashboard/policies/',
        ):
            self.assertNotContains(response, 'href="' + removed_path + '"')

    def test_signups_page_shows_graph_total_and_per_user_amount(self):
        response = self.client.get('/store/dashboard/signups/')

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.context['total_amount_paid'], Decimal('1250.50'))
        self.assertEqual(len(response.context['signup_chart']), 7)
        self.assertEqual(response.context['signups_today'], 2)
        self.assertEqual(response.context['signups_last_7_days'], 2)
        self.assertEqual(response.context['signups_this_month'], 2)
        self.assertEqual(response.context['paid_users'], 1)
        self.assertEqual(response.context['no_recorded_payment'], 1)
        self.assertEqual(response.context['paid_users_pct'], 50)
        self.assertEqual(
            [(role['label'], role['count']) for role in response.context['account_roles']],
            [('Customers', 1), ('Staff', 0), ('Superusers', 1)],
        )
        self.assertContains(response, 'Signups — last 7 days')
        self.assertContains(response, 'Total Amount Paid')
        self.assertContains(response, 'Joined Today')
        self.assertContains(response, 'Recorded payment coverage')
        self.assertContains(response, 'Account roles')
        self.assertContains(response, 'Location Enabled')
        self.assertContains(response, '20.2961, 85.8245')
        self.assertContains(response, '1250.50')
        self.assertContains(response, 'name="amount_paid"')

    def test_manual_user_amount_paid_is_optional_and_saved_when_provided(self):
        response = self.client.post('/store/dashboard/users/add/', {
            'next': 'dashboard_signups', 'name': 'Paid Customer',
            'email': 'paid-customer@example.com', 'phone': '9222222222',
            'password': 'customer-password', 'amount_paid': '499.99',
        })

        self.assertRedirects(response, '/store/dashboard/signups/')
        paid_profile = StoreProfile.objects.get(user__email='paid-customer@example.com')
        self.assertEqual(paid_profile.manual_amount_paid, Decimal('499.99'))

        optional_response = self.client.post('/store/dashboard/users/add/', {
            'next': 'dashboard_signups', 'name': 'Free Customer',
            'email': 'free-customer@example.com', 'phone': '9333333333',
            'password': 'customer-password', 'amount_paid': '',
        })
        self.assertRedirects(optional_response, '/store/dashboard/signups/')
        free_profile = StoreProfile.objects.get(user__email='free-customer@example.com')
        self.assertEqual(free_profile.manual_amount_paid, Decimal('0.00'))


class PWAFrontendSettingsTests(TestCase):
    def setUp(self):
        self.media_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.media_dir.cleanup)
        media_override = override_settings(MEDIA_ROOT=self.media_dir.name)
        media_override.enable()
        self.addCleanup(media_override.disable)

        self.admin = User.objects.create_superuser(
            username='pwa-admin@example.com', email='pwa-admin@example.com',
            password='admin-password',
        )
        StoreProfile.objects.create(user=self.admin, phone='9111111111')
        self.client.force_login(self.admin)

    def _icon_upload(self):
        image_bytes = io.BytesIO()
        Image.new('RGB', (700, 500), (12, 140, 90)).save(image_bytes, format='PNG')
        return SimpleUploadedFile('custom-pwa.png', image_bytes.getvalue(), content_type='image/png')

    def _save_enabled_settings(self):
        return self.client.post('/store/dashboard/pwa-settings/', {
            'is_enabled': 'on',
            'app_name': 'Rudra Custom AI',
            'short_name': 'Rudra AI',
            'description': 'Custom install description from the dashboard.',
            'theme_color': '#123456',
            'background_color': '#fedcba',
            'icon': self._icon_upload(),
        })

    def test_admin_pwa_settings_appear_on_homepage_and_manifest(self):
        saved = self._save_enabled_settings()
        self.assertEqual(saved.status_code, 200)
        self.assertTrue(saved.context['saved'])

        homepage = self.client.get('/')
        self.assertContains(homepage, 'Install Rudra Custom AI as an app')
        self.assertContains(homepage, '<meta name="theme-color" content="#123456">', html=True)
        self.assertContains(homepage, '/AI/manifest.json?v=')
        self.assertContains(homepage, '/AI/pwa-icon/192.png?v=')
        self.assertContains(homepage, "navigator.serviceWorker.register('/sw.js', { scope: '/' })")

        manifest_response = self.client.get('/AI/manifest.json')
        manifest = manifest_response.json()
        self.assertEqual(manifest['name'], 'Rudra Custom AI')
        self.assertEqual(manifest['short_name'], 'Rudra AI')
        self.assertEqual(manifest['description'], 'Custom install description from the dashboard.')
        self.assertEqual(manifest['theme_color'], '#123456')
        self.assertEqual(manifest['background_color'], '#fedcba')
        self.assertEqual(manifest['start_url'], '/')
        self.assertEqual(manifest['scope'], '/')
        self.assertIn('/AI/pwa-icon/512.png?v=', manifest['icons'][1]['src'])
        self.assertIn('no-store', manifest_response['Cache-Control'])

    def test_uploaded_icon_is_rendered_at_real_manifest_dimensions(self):
        self._save_enabled_settings()

        for size in (192, 512):
            with self.subTest(size=size):
                response = self.client.get(f'/AI/pwa-icon/{size}.png')
                self.assertEqual(response.status_code, 200)
                self.assertEqual(response['Content-Type'], 'image/png')
                rendered = Image.open(io.BytesIO(response.content))
                self.assertEqual(rendered.size, (size, size))

    def test_disabling_pwa_removes_manifest_banner_and_registration(self):
        self._save_enabled_settings()
        disabled = self.client.post('/store/dashboard/pwa-settings/', {
            'app_name': 'Rudra Custom AI',
            'short_name': 'Rudra AI',
            'description': 'Custom install description from the dashboard.',
            'theme_color': '#123456',
            'background_color': '#fedcba',
        })
        self.assertEqual(disabled.status_code, 200)

        homepage = self.client.get('/')
        self.assertNotContains(homepage, 'rel="manifest"')
        self.assertNotContains(homepage, 'id="installBanner"')
        self.assertNotContains(homepage, "navigator.serviceWorker.register('/sw.js'")
        self.assertContains(homepage, 'navigator.serviceWorker.getRegistrations()')


class DashboardBackupDeletionTests(TestCase):
    def setUp(self):
        self.admin = User.objects.create_superuser(
            username='backup-admin@example.com', email='backup-admin@example.com', password='admin-password',
        )
        StoreProfile.objects.create(user=self.admin, phone='9444444444')
        self.client.force_login(self.admin)

    def test_backup_page_contains_typed_delete_all_confirmation(self):
        response = self.client.get('/store/dashboard/backup/')

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, 'Delete All Backups')
        self.assertContains(response, 'name="confirmation"')
        self.assertContains(response, 'pattern="DELETE ALL"')
        self.assertContains(response, '/store/dashboard/backup/delete-all/')

    @patch('myapp.views.SiteCustomization.get_solo', side_effect=OperationalError('no such table'))
    def test_pages_survive_an_old_backup_before_customization_migration(self, get_solo):
        context = site_customization_context(RequestFactory().get('/store/'))

        self.assertEqual(context, {'SITE_FAVICON_URL': None})

    @patch('myapp.views.call_command')
    @patch('myapp.views.dropbox_backup.restore_backup')
    def test_restore_automatically_migrates_an_older_backup(self, restore_backup, call_command):
        with patch('django.contrib.sessions.backends.db.SessionStore.save') as session_save:
            response = self.client.post(
                '/store/dashboard/backup/restore/', {'filename': 'older-backup.sqlite3'},
            )

        self.assertRedirects(response, '/store/dashboard/backup/', fetch_redirect_response=False)
        restore_backup.assert_called_once()
        call_command.assert_called_once_with('migrate', interactive=False, verbosity=0)
        self.assertTrue(any(call.kwargs.get('must_create') for call in session_save.call_args_list))

    @patch('myapp.views.dropbox_backup.delete_all_backups')
    def test_delete_all_view_requires_exact_confirmation(self, delete_all):
        response = self.client.post(
            '/store/dashboard/backup/delete-all/', {'confirmation': 'delete all'}, follow=True,
        )

        self.assertEqual(response.status_code, 200)
        delete_all.assert_not_called()
        self.assertContains(response, 'Type DELETE ALL exactly')

    @patch('myapp.views.dropbox_backup.delete_all_backups', return_value=3)
    def test_delete_all_view_calls_dropbox_after_confirmation(self, delete_all):
        response = self.client.post(
            '/store/dashboard/backup/delete-all/', {'confirmation': 'DELETE ALL'}, follow=True,
        )

        self.assertEqual(response.status_code, 200)
        delete_all.assert_called_once()
        self.assertContains(response, 'Deleted 3 Dropbox backup items.')

    def test_helper_deletes_only_entries_inside_backup_folder_including_latest(self):
        dbx = Mock()
        dbx.files_list_folder.return_value = SimpleNamespace(
            entries=[
                SimpleNamespace(path_lower='/edutrellis store/backups/db_20260831.sqlite3'),
                SimpleNamespace(path_lower='/edutrellis store/backups/db_latest.sqlite3'),
                SimpleNamespace(path_lower='/unrelated/never-delete.sqlite3'),
            ],
            has_more=False,
        )

        with patch('myapp.dropbox_backup._client', return_value=dbx):
            deleted = dropbox_backup.delete_all_backups(SimpleNamespace())

        self.assertEqual(deleted, 2)
        self.assertEqual(dbx.files_delete_v2.call_count, 2)
        dbx.files_delete_v2.assert_any_call('/edutrellis store/backups/db_20260831.sqlite3')
        dbx.files_delete_v2.assert_any_call('/edutrellis store/backups/db_latest.sqlite3')


class RemovedPublicSurfaceTests(TestCase):
    def test_storefront_and_websitecreation_urls_are_gone(self):
        for path in (
            '/store/', '/store', '/estore', '/estore/',
            '/websitecreation/', '/websitecreation/contact/',
            '/contact/', '/store/api/cart/', '/store/product/anything/',
            '/store/policy/privacy/',
        ):
            self.assertEqual(self.client.get(path).status_code, 404, msg=path)

    def test_404_uses_the_saved_ai_frontend_theme(self):
        response = self.client.get('/missing-page/')
        self.assertContains(response, "localStorage.getItem('ai_theme')", status_code=404)
        self.assertContains(response, "localStorage.getItem('ai_color_theme')", status_code=404)
        self.assertContains(response, 'data-accent="blue"', status_code=404)
        self.assertContains(response, '--red:#059669', status_code=404)

    def test_emerald_default_is_migrated_for_every_existing_device(self):
        for path in ('/AI/', '/missing-page/'):
            response = self.client.get(path)
            self.assertContains(response, "var colorDefaultVersion = 'emerald-v1'", status_code=response.status_code)
            self.assertContains(response, "localStorage.setItem('ai_color_theme', 'emerald')", status_code=response.status_code)

        staff = User.objects.create_user('emerald-admin', password='password', is_staff=True)
        StoreProfile.objects.create(user=staff)
        self.client.force_login(staff)
        response = self.client.get('/store/dashboard/')
        self.assertContains(response, "var colorDefaultVersion = 'emerald-v1'")
        self.assertContains(response, "localStorage.setItem('ai_color_theme', 'emerald')")

    def test_ai_and_dashboard_routes_remain(self):
        self.assertEqual(self.client.get('/').status_code, 200)
        self.assertEqual(self.client.get('/AI/').status_code, 200)

        staff = User.objects.create_user('dashboard-admin', password='password', is_staff=True)
        StoreProfile.objects.create(user=staff)
        self.client.force_login(staff)
        self.assertEqual(self.client.get('/store/dashboard/').status_code, 200)

    def test_mobile_feature_intro_lists_images_and_models_and_closes_model_stack(self):
        response = self.client.get('/AI/')
        # Deliberately "Generate images", not "Generate and edit images":
        # editing an uploaded photo is unavailable upstream, so advertising it
        # here sent users straight into an error (AIReports #37, #44, #46).
        self.assertContains(response, 'Generate images')
        self.assertNotContains(response, 'Generate and edit images')
        self.assertContains(response, 'Access multiple AI models')
        self.assertContains(response, "localStorage.setItem('ai_model_intro_seen', '1')")
        self.assertContains(response, 'if (modelDropdown) modelDropdown.hidden = true')
        self.assertContains(response, 'max-height:calc(100dvh - 24px)')

    def test_homepage_shows_no_starter_questions_or_model_hint(self):
        response = self.client.get('/AI/')
        for removed in (
            'Help me pick a service',
            'How to use ChatGPT 5.6',
            'What is Vidhyora?',
            'Write something creative',
            'Explain something simply',
            'is ready — just type below',
            'from the model box below',
        ):
            with self.subTest(removed=removed):
                self.assertNotContains(response, removed)

    def test_signed_out_visitor_sees_the_neutral_description(self):
        response = self.client.get('/AI/')
        self.assertContains(response, 'A powerful AI with access to multiple models')
        # The rendered element, not the class name — '.hl-model' also appears
        # in the stylesheet, which is served to everyone.
        self.assertNotContains(response, '<span class="hl-model">')

    def test_signed_in_user_sees_the_chatgpt_description(self):
        staff = User.objects.create_user('full-access', password='password', is_staff=True)
        StoreProfile.objects.create(user=staff)
        self.client.force_login(staff)

        response = self.client.get('/AI/')
        self.assertContains(response, '<span class="hl-model">')
        self.assertContains(response, 'Full access to')
        self.assertNotContains(response, 'A powerful AI with access to multiple models')

    def test_apex_domain_redirects_to_ai_homepage(self):
        middleware = CanonicalHostMiddleware(lambda request: HttpResponse('page'))
        response = middleware(RequestFactory().get('/', HTTP_HOST='edutrellis.in', secure=True))
        self.assertEqual(response.status_code, 301)
        self.assertEqual(response['Location'], 'https://www.edutrellis.in/')

    def test_ai_assets_receive_browser_cache_headers(self):
        middleware = PublicAssetCacheMiddleware(lambda request: HttpResponse('asset'))
        response = middleware(RequestFactory().get('/static/ai-icon-192.png'))
        self.assertIn('max-age=86400', response['Cache-Control'])


class HTMLExtractionTests(TestCase):
    HTML = b'''<!doctype html>
        <html><head><style>.hidden { display:none }</style></head>
        <body><h1>Upload title</h1><p>Tom &amp; Jerry</p>
        <script>stealSecret()</script><noscript>hidden fallback</noscript></body></html>'''

    def test_html_uses_standard_library_fallback_without_bs4(self):
        with patch.dict('sys.modules', {'bs4': None}):
            text, truncated = doc_extract.extract('example.html', self.HTML)

        self.assertIn('Upload title', text)
        self.assertIn('Tom & Jerry', text)
        self.assertNotIn('stealSecret', text)
        self.assertNotIn('display:none', text)
        self.assertNotIn('hidden fallback', text)
        self.assertFalse(truncated)

    def test_html_upload_endpoint_returns_extracted_text(self):
        upload = SimpleUploadedFile('example.html', self.HTML, content_type='text/html')

        response = self.client.post('/AI/api/extract/', {'file': upload})

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()['status'], 'ok')
        self.assertIn('Upload title', response.json()['text'])
        self.assertIn('<h1>Upload title</h1>', response.json()['coding_text'])
        self.assertIn('<script>stealSecret()</script>', response.json()['coding_text'])

    def test_common_source_code_file_is_supported_without_renaming(self):
        source = b'def greet(name):\n    return f"Hello {name}"\n'

        text, truncated = doc_extract.extract('app.py', source)
        coding_text, coding_truncated = doc_extract.extract_editable_source('app.py', source, text)

        self.assertEqual(text, source.decode())
        self.assertEqual(coding_text, source.decode())
        self.assertFalse(truncated)
        self.assertFalse(coding_truncated)

    def test_document_action_instructions_keep_coding_and_details_separate(self):
        coding = _ai_document_instruction('coding', 'index.html')
        details = _ai_document_instruction('details', 'index.html')

        self.assertIn('COMPLETE updated file', coding)
        self.assertIn('never return only a patch', coding)
        self.assertIn('Analyse and explain only', details)
        self.assertIn('Do not rewrite the file', details)


class AIDocumentActionTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(
            username='document-actions@example.com',
            email='document-actions@example.com',
            password='test-password-123',
            is_staff=True,
        )
        self.client.force_login(self.user)

    def test_coding_action_forces_code_mode_and_full_file_instruction(self):
        payload = {
            'message': 'Change the theme colors to blue.',
            'model': 'ultra',
            'document_name': 'index.html',
            'document_text': '<html><body>Original</body></html>',
            'document_mode': 'coding',
            'document_truncated': False,
        }
        with patch('myapp.views.ai_chat.stream_chat', return_value=iter(['updated file'])) as stream_chat:
            response = self.client.post(
                '/AI/api/send/', data=json.dumps(payload), content_type='application/json'
            )
            body = b''.join(response.streaming_content).decode()

        self.assertEqual(response.status_code, 200)
        self.assertEqual(body, 'updated file')
        kwargs = stream_chat.call_args.kwargs
        self.assertEqual(kwargs['model_key'], 'code')
        self.assertEqual(kwargs['max_tokens'], 6000)
        self.assertIn('COMPLETE updated file', kwargs['document_instruction'])

    def test_details_action_keeps_analysis_only_instruction(self):
        payload = {
            'message': 'Show details about this file only.',
            'model': 'quick',
            'document_name': 'report.pdf',
            'document_text': 'Quarterly report content',
            'document_mode': 'details',
        }
        with patch('myapp.views.ai_chat.stream_chat', return_value=iter(['details'])) as stream_chat:
            response = self.client.post(
                '/AI/api/send/', data=json.dumps(payload), content_type='application/json'
            )
            b''.join(response.streaming_content)

        kwargs = stream_chat.call_args.kwargs
        self.assertEqual(kwargs['model_key'], 'quick')
        self.assertIsNone(kwargs['max_tokens'])
        self.assertIn('Do not rewrite the file', kwargs['document_instruction'])
